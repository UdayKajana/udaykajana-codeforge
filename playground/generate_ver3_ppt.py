"""
DSPy Evaluation v3 — Visual-first PPT
  Visual cards, score bars, spec badges, timeline, minimal text.
  Run:  python generate_ver3_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Palette ──────────────────────────────────────────────────────────────────
P_RED      = RGBColor(0xEB, 0x26, 0x2A)
P_DARK     = RGBColor(0x1C, 0x1C, 0x1C)
P_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
P_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
P_MGRAY    = RGBColor(0xAB, 0xAB, 0xAB)
P_DGRAY    = RGBColor(0x4A, 0x4A, 0x4A)
P_GREEN    = RGBColor(0x1B, 0x87, 0x3E)
P_GREEN_LT = RGBColor(0xE6, 0xF7, 0xEC)
P_RED_LT   = RGBColor(0xFD, 0xE8, 0xE8)
P_AMBER    = RGBColor(0xB3, 0x5A, 0x00)
P_AMBER_LT = RGBColor(0xFE, 0xF3, 0xE0)
P_BORDER   = RGBColor(0xDC, 0xDC, 0xDC)
P_BLUE     = RGBColor(0x18, 0x65, 0xD6)
P_BLUE_LT  = RGBColor(0xE8, 0xF0, 0xFE)
P_TEAL     = RGBColor(0x00, 0x82, 0x8A)
P_TEAL_LT  = RGBColor(0xE0, 0xF5, 0xF7)
P_PURPLE   = RGBColor(0x6B, 0x2F, 0x9E)
P_PURP_LT  = RGBColor(0xF0, 0xE8, 0xFA)
P_CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
P_SLATE    = RGBColor(0x3A, 0x3F, 0x4A)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation("pod_copy- Copy.pptx")
xml_slides = prs.slides._sldIdLst
while len(xml_slides):
    rId = xml_slides[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[0])

COVER = prs.slide_layouts[1]
BLANK = prs.slide_layouts[42]


# ── Core helpers ─────────────────────────────────────────────────────────────
def bg(slide, color=P_WHITE):
    b = slide.background; b.fill.solid(); b.fill.fore_color.rgb = color

def box(sl, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def tb(sl, text, x, y, w, h, sz=Pt(10), bold=False,
       color=P_DARK, align=PP_ALIGN.LEFT, italic=False):
    t = sl.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return t


def symbol_tb(sl, symbol, text, x, y, w, h, sym_color, sym_sz=Pt(9), text_sz=Pt(8)):
    """Colored symbol run + grey text run in one textbox."""
    t = sl.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = symbol + "  "
    r1.font.size = sym_sz; r1.font.bold = True
    r1.font.color.rgb = sym_color
    r2 = p.add_run()
    r2.text = text
    r2.font.size = text_sz; r2.font.bold = False
    r2.font.color.rgb = P_DGRAY
    return t
def hdr(sl, title, subtitle=""):
    box(sl, 0, 0, W, Inches(1.0), P_RED)
    tb(sl, title, Inches(0.45), Inches(0.10), Inches(12.0), Inches(0.54),
       Pt(20), True, P_WHITE)
    if subtitle:
        tb(sl, subtitle, Inches(0.45), Inches(0.64), Inches(12.0), Inches(0.32),
           Pt(9.5), False, RGBColor(0xFF, 0xC0, 0xBF))

def pill(sl, label, x, y, w, h, bg_col, fg_col):
    """Rounded-look badge / pill."""
    box(sl, x, y, w, h, bg_col, fg_col, Pt(0.75))
    tb(sl, label, x + Inches(0.06), y + Inches(0.04),
       w - Inches(0.1), h - Inches(0.06), Pt(8), True, fg_col, PP_ALIGN.CENTER)

def score_bar(sl, x, y, w, h, value_0_1, color):
    """Visual progress bar: grey track + colored fill."""
    box(sl, x, y, w, h, RGBColor(0xE0, 0xE0, 0xE0))
    fill = max(0.0, min(1.0, value_0_1))
    if fill > 0:
        box(sl, x, y, w * fill, h, color)

def section_label(sl, label, x, y, w, fg=P_MGRAY):
    tb(sl, label.upper(), x, y, w, Inches(0.2), Pt(7.5), True, fg)


# =============================================================================
# SLIDE 1 — Cover
# =============================================================================
def slide_cover():
    sl = prs.slides.add_slide(COVER)
    for ph in sl.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "DSPy vs No-DSPy"
            ph.text_frame.paragraphs[0].runs[0].font.bold = True
        elif idx == 1:
            ph.text = "A Comparative Evaluation"


# =============================================================================
# SLIDE 2 — Visual Email (left) + Flow + DSPy Intro (right)
# =============================================================================
def slide_flow():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    hdr(sl,
        "Pipeline Architecture: 2-Agent Direct vs DSPy Orchestrator",
        "Real Japanese procurement email (Suzuki Sekkei RFQ) — two pipelines, one email, very different outputs")

    CY = Inches(1.08)
    CH = Inches(6.28)    # content height

    # ============================================================
    # LEFT PANEL — Visual Email Card (0.22" → 5.65")
    # ============================================================
    EX = Inches(0.22)
    EW = Inches(5.44)

    box(sl, EX, CY, EW, CH, P_WHITE, P_BORDER)

    # Header strip
    box(sl, EX, CY, EW, Inches(0.34), P_CHARCOAL)
    tb(sl, "TEST EMAIL  —  Suzuki Sekkei RFQ  (Japanese)", EX + Inches(0.12),
       CY + Inches(0.05), EW - Inches(0.2), Inches(0.26), Pt(9), True, P_WHITE)

    cy = CY + Inches(0.44)   # running cursor

    # ── Sender / Header fields ────────────────────────────────
    for label, val, vc in [
        ("FROM",    "Kenta Yamamoto  [山本健太]",        P_DARK),
        ("",        "yamamoto.kenta@suzuki-sekkei.co.jp", P_MGRAY),
        ("TO",      "bids@globalparts.com",               P_DARK),
        ("DATE",    "May 26, 2026   10:32 AM  JST",       P_DARK),
        ("SUBJECT", "[URGENT] Procurement — 6-Axis Arc Welding Robots", P_RED),
    ]:
        if label:
            tb(sl, label, EX + Inches(0.14), cy, Inches(0.72), Inches(0.22),
               Pt(7.5), True, P_MGRAY)
            tb(sl, val, EX + Inches(0.9), cy, EW - Inches(1.0), Inches(0.22),
               Pt(8), label == "SUBJECT", vc)
        else:
            tb(sl, val, EX + Inches(0.9), cy, EW - Inches(1.0), Inches(0.22),
               Pt(7.5), False, P_MGRAY, PP_ALIGN.LEFT, True)
        cy += Inches(0.24)

    # Divider
    box(sl, EX + Inches(0.14), cy, EW - Inches(0.28), Inches(0.02), P_BORDER)
    cy += Inches(0.1)

    # ── Item Section ─────────────────────────────────────────
    section_label(sl, "Procurement Item", EX + Inches(0.14), cy, EW - Inches(0.28))
    cy += Inches(0.22)

    tb(sl, "6-Axis Arc Welding Robot",
       EX + Inches(0.14), cy, EW - Inches(0.28), Inches(0.28), Pt(12), True, P_DARK)
    cy += Inches(0.30)

    # Model pills (2 small horizontal pills)
    pill(sl, "FANUC M-10iA/12",     EX + Inches(0.14), cy, Inches(1.9), Inches(0.28),
         P_BLUE_LT, P_BLUE)
    pill(sl, "Yaskawa MA1440",       EX + Inches(2.1),  cy, Inches(1.5), Inches(0.28),
         P_TEAL_LT, P_TEAL)
    pill(sl, "Equiv. OK",            EX + Inches(3.66), cy, Inches(0.9), Inches(0.28),
         P_LGRAY, P_MGRAY)
    cy += Inches(0.36)

    # Qty + Budget row
    box(sl, EX + Inches(0.14), cy, Inches(1.6), Inches(0.55), P_LGRAY, P_BORDER, Pt(0.5))
    tb(sl, "QTY", EX + Inches(0.24), cy + Inches(0.04), Inches(1.4), Inches(0.18),
       Pt(7), True, P_MGRAY)
    tb(sl, "18 units", EX + Inches(0.24), cy + Inches(0.22), Inches(1.4), Inches(0.26),
       Pt(13), True, P_DARK, PP_ALIGN.LEFT)

    box(sl, EX + Inches(1.86), cy, Inches(3.4), Inches(0.55), P_DARK, None)
    tb(sl, "BUDGET", EX + Inches(1.98), cy + Inches(0.04), Inches(3.2), Inches(0.18),
       Pt(7), True, P_MGRAY)
    tb(sl, "JPY 220,000,000", EX + Inches(1.98), cy + Inches(0.22), Inches(3.2), Inches(0.26),
       Pt(13), True, P_WHITE)
    cy += Inches(0.65)

    # Budget note
    tb(sl, "Excl. tax  •  Delivery & install incl.  •  AMC billed separately",
       EX + Inches(0.14), cy, EW - Inches(0.28), Inches(0.20),
       Pt(7.5), False, P_MGRAY, PP_ALIGN.LEFT, True)
    cy += Inches(0.28)

    # ── Delivery Timeline ─────────────────────────────────────
    box(sl, EX + Inches(0.14), cy, EW - Inches(0.28), Inches(0.02), P_BORDER)
    cy += Inches(0.08)
    section_label(sl, "Delivery Schedule  (Split: 3 × 6 units)", EX + Inches(0.14), cy, EW - Inches(0.28))
    cy += Inches(0.22)

    phases = [
        ("Phase 1", "Sep 15\n2026", P_TEAL,    P_TEAL_LT),
        ("Phase 2", "Nov 30\n2026", P_BLUE,    P_BLUE_LT),
        ("Phase 3", "Jan 31\n2027", P_PURPLE,  P_PURP_LT),
        ("Full Ops","Mar 1\n2027",  P_GREEN,   P_GREEN_LT),
    ]
    PH_W = Inches(1.05)
    ARR_W = Inches(0.18)
    total_w = len(phases) * PH_W + (len(phases) - 1) * ARR_W
    ph_start = EX + Inches(0.14) + (EW - Inches(0.28) - total_w) / 2
    PHASE_H = Inches(0.70)
    for i, (ph_lbl, ph_date, fc, bc) in enumerate(phases):
        px = ph_start + i * (PH_W + ARR_W)
        box(sl, px, cy, PH_W, PHASE_H, bc, fc, Pt(1.0))
        tb(sl, ph_lbl, px, cy + Inches(0.06), PH_W, Inches(0.20),
           Pt(7.5), True, fc, PP_ALIGN.CENTER)
        tb(sl, "6 units" if i < 3 else "All ops",
           px, cy + Inches(0.25), PH_W, Inches(0.18),
           Pt(7), False, P_DGRAY, PP_ALIGN.CENTER)
        tb(sl, ph_date, px, cy + Inches(0.44), PH_W, Inches(0.24),
           Pt(7.5), True, fc, PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            tb(sl, "→", px + PH_W, cy + Inches(0.25), ARR_W, Inches(0.22),
               Pt(11), True, P_MGRAY, PP_ALIGN.CENTER)
    cy += PHASE_H + Inches(0.14)

    # ── Tech Specs Badges ─────────────────────────────────────
    box(sl, EX + Inches(0.14), cy, EW - Inches(0.28), Inches(0.02), P_BORDER)
    cy += Inches(0.08)
    section_label(sl, "Technical Specifications", EX + Inches(0.14), cy, EW - Inches(0.28))
    cy += Inches(0.22)

    specs = [
        ("IP67",         P_BLUE_LT,   P_BLUE),
        ("±0.05 mm",     P_TEAL_LT,   P_TEAL),
        ("10 kg payload",P_GREEN_LT,  P_GREEN),
        ("1,400 mm reach",P_PURP_LT,  P_PURPLE),
        ("EtherNet/IP",  P_AMBER_LT,  P_AMBER),
    ]
    sx = EX + Inches(0.14)
    for sp_lbl, sp_bg, sp_fg in specs:
        sw = Inches(max(0.72, 0.12 + 0.068 * len(sp_lbl)))  # ~0.068" per char at Pt(7.5)
        if sx + sw > EX + EW - Inches(0.14):
            sx = EX + Inches(0.14); cy += Inches(0.30)
        box(sl, sx, cy, sw, Inches(0.26), sp_bg, sp_fg, Pt(0.75))
        tb(sl, sp_lbl, sx + Inches(0.07), cy + Inches(0.04),
           sw - Inches(0.1), Inches(0.20), Pt(7.5), True, sp_fg, PP_ALIGN.CENTER)
        sx += sw + Inches(0.1)
    cy += Inches(0.36)

    # ── Terms ─────────────────────────────────────────────────
    box(sl, EX + Inches(0.14), cy, EW - Inches(0.28), Inches(0.02), P_BORDER)
    cy += Inches(0.08)
    terms = [
        ("AMC",     "Min. 3 years — scope UNDEFINED in email",    P_RED),
        ("BID BY",  "June 20, 2026   17:00 JST",                  P_DARK),
        ("CONTACT", "Yamamoto Kenta [山本健太]  —  Head of Procurement", P_DARK),
    ]
    for t_lbl, t_val, t_vc in terms:
        tb(sl, t_lbl, EX + Inches(0.14), cy, Inches(0.7), Inches(0.22),
           Pt(7.5), True, P_MGRAY)
        tb(sl, t_val, EX + Inches(0.88), cy, EW - Inches(1.0), Inches(0.22),
           Pt(8), False, t_vc)
        cy += Inches(0.24)

    # ============================================================
    # RIGHT PANEL — DSPy Intro + Pipeline Flow (6.0" → 13.1")
    # ============================================================
    RX = Inches(5.85)
    RW = W - RX - Inches(0.22)

    # ── DSPy Intro: 2×2 card grid ─────────────────────────────
    INTRO_CARDS = [
        ("Structured\nPipelines",   "Chain LLM calls as typed, schema-enforced modules", P_PURPLE, P_PURP_LT),
        ("Auditable\nHandoffs",     "Every step traceable via inspect_history()", P_BLUE, P_BLUE_LT),
        ("Constrained\nExtraction", "Signatures enforce output schema — stops drift", P_TEAL, P_TEAL_LT),
        ("Auto-\nOptimizable",      "BootstrapFewShot tunes prompts automatically",  P_GREEN, P_GREEN_LT),
    ]
    CARD_W = (RW - Inches(0.12)) / 2
    CARD_H = Inches(1.08)
    intro_label_y = CY
    tb(sl, "WHY DSPY?", RX, intro_label_y, RW, Inches(0.22), Pt(9), True, P_SLATE)
    intro_label_y += Inches(0.24)

    for i, (title, desc, fc, bc) in enumerate(INTRO_CARDS):
        col = i % 2; row = i // 2
        cx = RX + col * (CARD_W + Inches(0.12))
        cy_c = intro_label_y + row * (CARD_H + Inches(0.1))
        box(sl, cx, cy_c, CARD_W, CARD_H, bc, fc, Pt(1.0))
        box(sl, cx, cy_c, Inches(0.08), CARD_H, fc)
        tb(sl, title, cx + Inches(0.16), cy_c + Inches(0.08),
           CARD_W - Inches(0.22), Inches(0.38), Pt(10), True, fc)
        tb(sl, desc, cx + Inches(0.16), cy_c + Inches(0.52),
           CARD_W - Inches(0.22), Inches(0.50), Pt(8), False, P_DGRAY)

    FLOW_START_Y = intro_label_y + 2 * CARD_H + Inches(0.3)

    # ── Pipeline Flow Diagram ─────────────────────────────────
    box(sl, RX, FLOW_START_Y, RW, CY + CH - FLOW_START_Y, P_LGRAY, P_BORDER)
    box(sl, RX, FLOW_START_Y, RW, Inches(0.28), P_CHARCOAL)
    tb(sl, "PIPELINE COMPARISON", RX + Inches(0.14), FLOW_START_Y + Inches(0.04),
       RW - Inches(0.2), Inches(0.22), Pt(8.5), True, P_WHITE)

    def draw_pipeline(base_y, case_label, label_color, nodes):
        tb(sl, case_label, RX + Inches(0.14), base_y, RW - Inches(0.2), Inches(0.20),
           Pt(8.5), True, label_color)
        BH  = Inches(0.60)
        BY  = base_y + Inches(0.24)
        AY  = BY + BH / 2 - Inches(0.1)
        x   = RX + Inches(0.14)
        for j, (n1, n2, nw, nfill, nborder) in enumerate(nodes):
            box(sl, x, BY, nw, BH, nfill, nborder, Pt(1.0))
            tb(sl, n1, x + Inches(0.04), BY + Inches(0.06),
               nw - Inches(0.08), Inches(0.22), Pt(8.5), True,
               P_DGRAY if nborder == P_BORDER else nborder, PP_ALIGN.CENTER)
            if n2:
                tb(sl, n2, x + Inches(0.04), BY + Inches(0.30),
                   nw - Inches(0.08), Inches(0.26), Pt(7), False, P_DGRAY, PP_ALIGN.CENTER)
            x += nw
            if j < len(nodes) - 1:
                tb(sl, "→", x, AY, Inches(0.22), Inches(0.26),
                   Pt(12), True, P_MGRAY, PP_ALIGN.CENTER)
                x += Inches(0.22)
        return BY + BH

    # Case 1 - No DSPy
    r1_y   = FLOW_START_Y + Inches(0.36)
    r1_end = draw_pipeline(r1_y, "Case 1 — No DSPy (2-Agent Direct)", P_RED, [
        ("Email",   None,             Inches(0.8),  P_LGRAY,   P_BORDER),
        ("Agent-1", "Extract+Audit",  Inches(2.0),  P_RED_LT,  P_RED),
        ("Agent-2", "Enrich+Validate",Inches(2.0),  P_RED_LT,  P_RED),
        ("Output",  "JSON",           Inches(0.8),  P_LGRAY,   P_BORDER),
    ])

    # Case 2 - DSPy
    r2_y   = r1_end + Inches(0.32)
    r2_end = draw_pipeline(r2_y, "Case 2 — With DSPy (Orchestrator)", P_GREEN, [
        ("Email",      None,           Inches(0.8),  P_LGRAY,   P_BORDER),
        ("DSPy",       "Predict(Extr)",Inches(1.5),  P_PURP_LT, P_PURPLE),
        ("DSPy",       "Predict(Audit)",Inches(1.5), P_PURP_LT, P_PURPLE),
        ("Agent",      "Enrich",       Inches(1.2),  P_GREEN_LT,P_GREEN),
        ("Output",     "JSON",         Inches(0.8),  P_LGRAY,   P_BORDER),
    ])

    # Key callout box
    CALL_Y = r2_end + Inches(0.16)
    box(sl, RX + Inches(0.14), CALL_Y, RW - Inches(0.28), Inches(0.44), P_AMBER_LT, P_AMBER)
    tb(sl, "Key Difference:  DSPy typed Signatures enforce schema at each step — "
          "prevents phantom fields and missing domain defaults",
       RX + Inches(0.26), CALL_Y + Inches(0.08), RW - Inches(0.52), Inches(0.30),
       Pt(8.5), False, P_AMBER)


# =============================================================================
# SHARED: Visual Metric Row
# =============================================================================
SCORE_W = Inches(1.85)

def metric_row(sl, T_X, T_Y, T_W, ROW_H,
               metric_name, score_text, bar_value, badge, badge_color,
               good_pts, bad_pts):
    """
    One visual metric row:
      [Score box | progress bar] [✓ Green col] [✗ Red col]
    good_pts / bad_pts: list of short strings (max ~40 chars each)
    """
    HALF_W = (T_W - SCORE_W) / 2

    # ── Score box ─────────────────────────────────────────────
    box(sl, T_X, T_Y, SCORE_W, ROW_H, badge_color[1], badge_color[0], Pt(1.5))

    # Big score number
    tb(sl, score_text,
       T_X + Inches(0.04), T_Y + Inches(0.10),
       SCORE_W - Inches(0.08), Inches(0.42),
       Pt(26), True, badge_color[0], PP_ALIGN.CENTER)

    # Metric name
    tb(sl, metric_name,
       T_X + Inches(0.04), T_Y + Inches(0.54),
       SCORE_W - Inches(0.08), Inches(0.22),
       Pt(8), True, P_DGRAY, PP_ALIGN.CENTER)

    # Badge
    bw = Inches(0.56)
    bx = T_X + (SCORE_W - bw) / 2
    box(sl, bx, T_Y + ROW_H - Inches(0.26), bw, Inches(0.22), badge_color[0])
    tb(sl, badge, bx, T_Y + ROW_H - Inches(0.26), bw, Inches(0.22),
       Pt(7.5), True, P_WHITE, PP_ALIGN.CENTER)

    # Progress bar strip
    BAR_H = Inches(0.10)
    bar_y = T_Y + ROW_H - Inches(0.50)
    score_bar(sl, T_X + Inches(0.14), bar_y,
              SCORE_W - Inches(0.28), BAR_H, bar_value, badge_color[0])

    # Thin divider
    box(sl, T_X + SCORE_W, T_Y, Inches(0.03), ROW_H, P_BORDER)

    # ── Green "What Worked" column ─────────────────────────────
    GX = T_X + SCORE_W + Inches(0.03)
    box(sl, GX, T_Y, HALF_W, Inches(0.24), P_GREEN_LT)
    box(sl, GX, T_Y, Inches(0.05), Inches(0.24), P_GREEN)
    tb(sl, "✓  What Worked",
       GX + Inches(0.1), T_Y + Inches(0.04), HALF_W - Inches(0.14), Inches(0.18),
       Pt(8), True, P_GREEN)

    for i, pt in enumerate(good_pts[:3]):
        by = T_Y + Inches(0.28) + i * Inches(0.23)
        symbol_tb(sl, "✓", pt, GX + Inches(0.1), by,
                  HALF_W - Inches(0.14), Inches(0.24), P_GREEN)

    # Thin divider
    box(sl, GX + HALF_W, T_Y, Inches(0.03), ROW_H, P_BORDER)

    # ── Red "What Failed" column ───────────────────────────────
    RDX = GX + HALF_W + Inches(0.03)
    box(sl, RDX, T_Y, HALF_W, Inches(0.24), P_RED_LT)
    box(sl, RDX, T_Y, Inches(0.05), Inches(0.24), P_RED)
    tb(sl, "✗  What Failed",
       RDX + Inches(0.1), T_Y + Inches(0.04), HALF_W - Inches(0.14), Inches(0.18),
       Pt(8), True, P_RED)

    for i, pt in enumerate(bad_pts[:3]):
        by = T_Y + Inches(0.28) + i * Inches(0.23)
        symbol_tb(sl, "✗", pt, RDX + Inches(0.1), by,
                  HALF_W - Inches(0.14), Inches(0.24), P_RED)

    # Full-row border
    box(sl, T_X, T_Y, T_W, ROW_H, RGBColor(0xFF,0xFF,0xFF), P_BORDER, Pt(0.3))


GOOD  = (P_GREEN,  P_GREEN_LT)
WARN  = (P_AMBER,  P_AMBER_LT)
POOR  = (P_RED,    P_RED_LT)


# =============================================================================
# SLIDE 3 — No-DSPy Visual Metrics
# =============================================================================
NODSPY_METRICS = [
    dict(
        metric_name="F1 Score",
        score_text="0.83",  bar_value=0.83,  badge="WARN",  badge_color=WARN,
        good_pts=["All 40 core facts extracted", "Specs, budget, delivery correct", "JIS / ISO standards identified"],
        bad_pts= ["2 phantom vendor sections", "12 NEEDS_INPUT — only 2 genuine", "Schema inflation lowers precision"],
    ),
    dict(
        metric_name="Faithfulness",
        score_text="0.95",  bar_value=0.95,  badge="GOOD",  badge_color=GOOD,
        good_pts=["Readable dates: 'Sep 15, 2026'", "Japanese brackets throughout", "'Equiv. models' captured"],
        bad_pts= ["Vendor fields not in source email", "Wrong risk flag for buyer RFQ", "Schema has phantom sections"],
    ),
    dict(
        metric_name="Output Relevancy",
        score_text="64%",   bar_value=0.64,  badge="POOR",  badge_color=POOR,
        good_pts=["Correct domain: procurement/RFQ", "All technical specs extracted", "open_items in array format"],
        bad_pts= ["14 NEEDS_INPUT (12 phantom)", "Zero domain defaults applied", "ERP record incomplete"],
    ),
    dict(
        metric_name="Output Correctness",
        score_text="0.80",  bar_value=0.80,  badge="WARN",  badge_color=WARN,
        good_pts=["Tech specs 100% correct", "Budget inclusions/exclusions right", "3-phase delivery accurate"],
        bad_pts= ["0/7 domain fields applied", "Only 45 of 55 fields filled", "Phantom fields inflate validation"],
    ),
    dict(
        metric_name="Hallucination",
        score_text="15%",   bar_value=0.15,  badge="POOR",  badge_color=POOR,
        good_pts=["Core facts grounded in email", "No numeric values fabricated", "Resolution log honest"],
        bad_pts= ["Vendor section from nothing", "12 phantom NEEDS_INPUT fields", "Wrong risk: 'no supplier details'"],
    ),
]

def slide_nodspy():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    hdr(sl, "No-DSPy Pipeline (2-Agent Direct) — Evaluation Metrics",
        "Pipeline B  ·  Japanese RFQ (Suzuki Sekkei)  ·  Two agents, no structured orchestration")

    T_X   = Inches(0.25)
    T_Y   = Inches(1.08)
    T_W   = W - Inches(0.50)
    # Available height: 7.5 - 1.08 (table top) - 0.08 (bottom gap) - 0.50 (summary) = 5.84"
    # Col header: 0.28", 5 rows
    COL_H = Inches(0.28)
    SUMM_H= Inches(0.50)
    ROW_H = (H - T_Y - Inches(0.08) - SUMM_H - COL_H) / 5  # ~1.09"

    # Column header strip
    HALF_W = (T_W - SCORE_W) / 2
    box(sl, T_X, T_Y, SCORE_W, COL_H, P_SLATE)
    tb(sl, "METRIC  /  SCORE", T_X + Inches(0.1), T_Y + Inches(0.05),
       SCORE_W - Inches(0.12), Inches(0.20), Pt(8), True, P_WHITE, PP_ALIGN.CENTER)
    GX2 = T_X + SCORE_W
    box(sl, GX2, T_Y, HALF_W, COL_H, P_GREEN)
    tb(sl, "WHAT WORKED  ✓", GX2 + Inches(0.1), T_Y + Inches(0.05),
       HALF_W - Inches(0.14), Inches(0.20), Pt(8), True, P_WHITE)
    NX2 = GX2 + HALF_W
    box(sl, NX2, T_Y, HALF_W, COL_H, P_RED)
    tb(sl, "WHAT FAILED  ✗", NX2 + Inches(0.1), T_Y + Inches(0.05),
       HALF_W - Inches(0.14), Inches(0.20), Pt(8), True, P_WHITE)

    for i, m in enumerate(NODSPY_METRICS):
        metric_row(sl, T_X, T_Y + COL_H + i * ROW_H, T_W, ROW_H, **m)

    # Summary banner
    SY = T_Y + COL_H + 5 * ROW_H + Inches(0.04)
    box(sl, T_X, SY, T_W, SUMM_H, P_RED_LT, P_RED)
    box(sl, T_X, SY, Inches(0.1), SUMM_H, P_RED)
    tb(sl, "Pipeline B  ·  Faithful to source (0.95) but schema inflation, missing domain fills,"
          " and wrong risk flags make it incomplete for ERP use. Loses 4 of 5 metrics.",
       T_X + Inches(0.2), SY + Inches(0.10), T_W - Inches(0.3), Inches(0.32),
       Pt(9), False, P_DARK)


# =============================================================================
# SLIDE 4 — DSPy Visual Metrics + Final Comparison Table
# =============================================================================
DSPY_METRICS = [
    dict(
        metric_name="F1 Score",
        score_text="0.87",  bar_value=0.87,  badge="GOOD",  badge_color=GOOD,
        good_pts=["55 fields (10 more than B)", "7 domain defaults auto-applied", "Only 2 genuine NEEDS_INPUT"],
        bad_pts= ["Duplicate keys in output", "open_items: {} not []", "amc_duration key duplicated"],
    ),
    dict(
        metric_name="Faithfulness",
        score_text="0.92",  bar_value=0.92,  badge="GOOD",  badge_color=GOOD,
        good_pts=["All 40 core facts verified", "Zero phantom sections", "Japanese fully translated"],
        bad_pts= ["ISO dates not readable English", "Company inferred from email", "3 risks absent from risk_flags"],
    ),
    dict(
        metric_name="Output Relevancy",
        score_text="79%",   bar_value=0.79,  badge="WARN",  badge_color=WARN,
        good_pts=["Full ERP-ready record", "Laws of Japan, JCAA, Net 30", "+23% vs Pipeline B"],
        bad_pts= ["risk_flags empty — 3 missed", "AMC scope, tax, RFQ ref gap", "2 NEEDS_INPUT remain"],
    ),
    dict(
        metric_name="Output Correctness",
        score_text="0.88",  bar_value=0.88,  badge="GOOD",  badge_color=GOOD,
        good_pts=["AuditSignature validates schema", "Domain fills auto-applied", "3-step structured pipeline"],
        bad_pts= ["open_items: {} schema error", "Dates in ISO not English", "Only 2 resolution entries"],
    ),
    dict(
        metric_name="Hallucination",
        score_text="8%",    bar_value=0.08,  badge="GOOD",  badge_color=GOOD,
        good_pts=["Zero phantom sections", "No vendor data inflated", "Halved vs Pipeline B"],
        bad_pts= ["Company from email address", "Silent miss on 3 real risks", "Shared risk blind spot"],
    ),
]

def slide_dspy_metrics():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    hdr(sl, "DSPy Orchestrator — Evaluation Metrics & Final Comparison",
        "Pipeline A  ·  Same Japanese RFQ  ·  Typed Signatures + structured 3-step audit pipeline")

    T_X = Inches(0.25)
    T_Y = Inches(1.08)
    T_W = W - Inches(0.50)

    # Summary table dimensions (pre-calculate to back into ROW_H)
    SUMM_HDR_H  = Inches(0.26)
    SUMM_COL_H  = Inches(0.28)
    SUMM_ROW_H  = Inches(0.28)
    SUMM_ROWS   = 5
    VERDICT_H   = Inches(0.34)
    GAP         = Inches(0.12)
    SUMM_TOTAL  = SUMM_HDR_H + SUMM_COL_H + SUMM_ROWS * SUMM_ROW_H + VERDICT_H
    COL_H       = Inches(0.28)

    available   = H - T_Y - Inches(0.06) - SUMM_TOTAL - GAP
    ROW_H       = (available - COL_H) / 5     # per-metric row

    # Column headers
    HALF_W = (T_W - SCORE_W) / 2
    box(sl, T_X, T_Y, SCORE_W, COL_H, P_SLATE)
    tb(sl, "METRIC  /  SCORE", T_X + Inches(0.1), T_Y + Inches(0.05),
       SCORE_W - Inches(0.12), Inches(0.20), Pt(8), True, P_WHITE, PP_ALIGN.CENTER)
    GX2 = T_X + SCORE_W
    box(sl, GX2, T_Y, HALF_W, COL_H, P_GREEN)
    tb(sl, "WHAT WORKED  ✓", GX2 + Inches(0.1), T_Y + Inches(0.05),
       HALF_W - Inches(0.14), Inches(0.20), Pt(8), True, P_WHITE)
    NX2 = GX2 + HALF_W
    box(sl, NX2, T_Y, HALF_W, COL_H, P_RED)
    tb(sl, "WHAT FAILED  ✗", NX2 + Inches(0.1), T_Y + Inches(0.05),
       HALF_W - Inches(0.14), Inches(0.20), Pt(8), True, P_WHITE)

    for i, m in enumerate(DSPY_METRICS):
        metric_row(sl, T_X, T_Y + COL_H + i * ROW_H, T_W, ROW_H, **m)

    # ── Final Summary Comparison Table ───────────────────────
    SY = T_Y + COL_H + 5 * ROW_H + GAP

    # Dark header
    box(sl, T_X, SY, T_W, SUMM_HDR_H, P_DARK)
    tb(sl, "FINAL COMPARISON  —  Pipeline A (DSPy)  vs  Pipeline B (2-Agent Direct)",
       T_X + Inches(0.14), SY + Inches(0.04), T_W - Inches(0.2), Inches(0.20),
       Pt(9), True, P_WHITE)

    # Summary columns
    SCW = [Inches(1.8), Inches(1.0), Inches(1.0), Inches(0.86)]
    SCW.append(T_W - sum(SCW))
    s_heads = ["Metric", "No-DSPy", "DSPy", "Delta", "Winner & Finding"]
    CHY = SY + SUMM_HDR_H
    x = T_X
    for i, (sh, cw) in enumerate(zip(s_heads, SCW)):
        box(sl, x, CHY, cw, SUMM_COL_H, P_CHARCOAL)
        align = PP_ALIGN.CENTER if 0 < i < 4 else PP_ALIGN.LEFT
        tb(sl, sh, x + Inches(0.06), CHY + Inches(0.06),
           cw - Inches(0.1), Inches(0.20), Pt(8), True, P_WHITE, align)
        x += cw

    s_rows = [
        ("F1 Score",         "0.83", "0.87", "+5%",   True,  "DSPy — domain defaults + fewer phantom fields"),
        ("Faithfulness",     "0.95", "0.92", "−3 pp", False, "No-DSPy — readable dates, Japanese brackets"),
        ("Output Relevancy", "64%",  "79%",  "+23%",  True,  "DSPy — ERP-ready, 55 fields, JCAA, Net 30"),
        ("Output Correct.",  "0.80", "0.88", "+10%",  True,  "DSPy — AuditSignature prevents schema drift"),
        ("Hallucination",    "15%",  "8%",   "−47%",  True,  "DSPy — zero phantom sections, halved rate"),
    ]
    for r, (mn, vb, va, delta, dspy_w, finding) in enumerate(s_rows):
        ry = CHY + SUMM_COL_H + r * SUMM_ROW_H
        rb = P_LGRAY if r % 2 == 0 else P_WHITE
        x  = T_X
        dc = P_GREEN if dspy_w else P_AMBER
        for c, (cell, cw) in enumerate(zip([mn, vb, va, delta, finding], SCW)):
            box(sl, x, ry, cw, SUMM_ROW_H, rb, P_BORDER, Pt(0.25))
            if c == 0:
                tb(sl, cell, x + Inches(0.06), ry + Inches(0.05),
                   cw - Inches(0.1), SUMM_ROW_H - Inches(0.06), Pt(8.5), True, P_DARK)
            elif c == 1:
                tb(sl, cell, x + Inches(0.04), ry + Inches(0.04),
                   cw - Inches(0.08), SUMM_ROW_H - Inches(0.06), Pt(10), True, P_RED, PP_ALIGN.CENTER)
            elif c == 2:
                tb(sl, cell, x + Inches(0.04), ry + Inches(0.04),
                   cw - Inches(0.08), SUMM_ROW_H - Inches(0.06), Pt(10), True,
                   P_GREEN if dspy_w else P_AMBER, PP_ALIGN.CENTER)
            elif c == 3:
                box(sl, x + Inches(0.04), ry + Inches(0.04),
                    cw - Inches(0.08), SUMM_ROW_H - Inches(0.08), dc)
                tb(sl, cell, x + Inches(0.04), ry + Inches(0.04),
                   cw - Inches(0.08), SUMM_ROW_H - Inches(0.08), Pt(8), True, P_WHITE, PP_ALIGN.CENTER)
            else:
                tb(sl, cell, x + Inches(0.08), ry + Inches(0.05),
                   cw - Inches(0.14), SUMM_ROW_H - Inches(0.06), Pt(8), False, P_DGRAY)
            x += cw

    # Verdict bar
    VY = CHY + SUMM_COL_H + SUMM_ROWS * SUMM_ROW_H + Inches(0.02)
    box(sl, T_X, VY, T_W, VERDICT_H, P_GREEN_LT, P_GREEN)
    box(sl, T_X, VY, Inches(0.1), VERDICT_H, P_GREEN)
    tb(sl, "VERDICT:  DSPy wins 4/5 — Hallucination 8% vs 15%  ·  Confidence 79% vs 64%  ·  "
          "55 vs 45 fields  ·  Zero phantom sections.   "
          "No-DSPy wins Faithfulness (0.95 vs 0.92).   "
          "Critical gap BOTH:  add RiskSignature for AMC, tax & RFQ ID.",
       T_X + Inches(0.2), VY + Inches(0.06), T_W - Inches(0.3), VERDICT_H - Inches(0.08),
       Pt(8.5), False, P_DARK)


# ── Build ─────────────────────────────────────────────────────────────────────
slide_cover()
slide_flow()
slide_nodspy()
slide_dspy_metrics()

OUT = "dspy_evaluation_ver3.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")

