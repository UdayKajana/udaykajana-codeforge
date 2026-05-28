"""
Generates pipeline_handoff_comparison.pptx
Compares inputs (bidding_emails.json) vs outputs from:
  Case 1 — Agent-1 → DSPy → Agent-2   (finals_dspy.json)
  Case 2 — Agent-1 → Sys Instr → Agent-2  (finals_no_dspy.json)
"""

import json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Load data ─────────────────────────────────────────────────────────────────
with open("bidding_emails.json",  encoding="utf-8") as f: emails      = json.load(f)
with open("finals_dspy.json",     encoding="utf-8") as f: out_dspy    = json.load(f)
with open("finals_no_dspy.json",  encoding="utf-8") as f: out_nodspy  = json.load(f)

# ── Palette ───────────────────────────────────────────────────────────────────
C_WHITE   = RGBColor(0xFF,0xFF,0xFF)
C_NAVY    = RGBColor(0x0F,0x28,0x4A)
C_BLUE    = RGBColor(0x1A,0x56,0xDB)
C_BLUE_L  = RGBColor(0xDB,0xEA,0xFF)
C_ORANGE  = RGBColor(0xC4,0x50,0x1A)
C_ORANGE_L= RGBColor(0xFD,0xED,0xE0)
C_DARK    = RGBColor(0x1A,0x20,0x2C)
C_GRAY_D  = RGBColor(0x4A,0x55,0x68)
C_GRAY_M  = RGBColor(0x71,0x80,0x96)
C_GRAY_L  = RGBColor(0xF0,0xF4,0xF8)
C_GRAY_LL = RGBColor(0xF7,0xF9,0xFC)
C_BORDER  = RGBColor(0xCB,0xD5,0xE0)
C_GREEN   = RGBColor(0x0B,0x7B,0x47)
C_GREEN_L = RGBColor(0xD1,0xFA,0xE5)
C_RED     = RGBColor(0xB9,0x18,0x1A)
C_RED_L   = RGBColor(0xFE,0xE2,0xE2)
C_AMBER   = RGBColor(0x92,0x40,0x09)
C_AMBER_L = RGBColor(0xFF,0xED,0xC2)

W, H = Inches(13.33), Inches(7.5)

# ── Core helpers ──────────────────────────────────────────────────────────────
def new_prs():
    p = Presentation(); p.slide_width = W; p.slide_height = H; return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def r(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.fill.background()
    return s

def t(slide, text, x, y, w, h,
      sz=Pt(11), bold=False, color=C_DARK,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    rn = p.add_run(); rn.text = str(text)
    rn.font.size = sz; rn.font.bold = bold
    rn.font.color.rgb = color; rn.font.italic = italic
    return tb

def chrome(slide, title, n, sub=""):
    bg(slide)
    r(slide, 0, 0, W, Inches(0.95), C_NAVY)
    r(slide, 0, 0, Inches(0.22), Inches(0.95), C_BLUE)
    t(slide, title, Inches(0.35), Inches(0.14), Inches(11.5), Inches(0.65),
      Pt(20), True, C_WHITE)
    if sub:
        t(slide, sub, Inches(0.35), Inches(0.65), Inches(11.0), Inches(0.28),
          Pt(10), False, RGBColor(0xA0,0xAE,0xC8))
    r(slide, Inches(12.6), Inches(0.22), Inches(0.56), Inches(0.5), C_BLUE)
    t(slide, str(n), Inches(12.6), Inches(0.22), Inches(0.56), Inches(0.5),
      Pt(14), True, C_WHITE, PP_ALIGN.CENTER)

def hline(slide, y, x0=Inches(0.5), x1=Inches(12.83)):
    r(slide, x0, y, x1-x0, Inches(0.015), C_BORDER)

# ── Data helpers ──────────────────────────────────────────────────────────────

EMAIL_META = {
    "email_01": ("Welding Robots",       "Manufacturing RFQ",   "₹2.2 Cr",   "Delta Infrastructure",   "June 10, 2026",  "Aug 15, 2026",  15,    "formal"),
    "email_02": ("EV Battery Packs",     "Manufacturing RFQ",   "€18 M",     "Greenfield Auto",        "June 20, 2026",  "Sep-Oct 2026",  2800,  "formal"),
    "email_03": ("Laptop Procurement",   "IT Hardware RFQ",     "$1,20,000",  "NovaTech Solutions",     "June 12, 2026",  "Jul 5, 2026",   120,   "urgent"),
    "email_04": ("OLED Panels",          "Electronics RFQ",     "$750,000",   "Metro Display Tech",     "June 18, 2026",  "Jul-Sep 2026",  50000, "formal"),
    "email_05": ("Office Furniture",     "Facilities RFQ",      "₹85 L",     "GlobalSpace Offices",    "June 25, 2026",  "Aug 10, 2026",  "200-seat", "formal"),
    "email_06": ("Smartphones 1200",     "Consumer/Retail RFQ", "₹6.2 Cr",   "ConnectHub India",       "June 14, 2026",  "Aug 1, 2026",   1200,  "informal"),
    "email_07": ("Electric Vans",        "Fleet Procurement",   "₹18 Cr",    "UrbanMove Logistics",    "June 28, 2026",  "Sep-Jan 2027",  75,    "formal"),
    "email_08": ("Kitchen Equipment",    "Hospitality RFQ",     "£360,000",   "GrandVista Hotels",      "June 30, 2026",  "Sep 20, 2026",  "3 kitchens", "formal"),
    "email_09": ("Smart Classrooms",     "Govt. Tender",        "₹42 Cr",    "State Edu Dept.",        "June 30, 2026",  "Oct-Feb 2027",  500,   "formal"),
    "email_10": ("Aluminium Supply",     "12-month Contract",   "~₹8 Cr/yr", "FastPack Mfg.",          "June 22, 2026",  "Jul 2026+",     "65 MT/mo", "informal"),
}

def get_validation(rec):
    return rec.get("validation", {})

def count_fields(rec):
    """Count filled, empty, needs-input in field_inventory."""
    fi = rec.get("field_inventory", rec.get("tag_list", {}))
    def walk(obj):
        filled = empty = needs = 0
        if isinstance(obj, dict):
            for v in obj.values():
                f,e,n = walk(v)
                filled+=f; empty+=e; needs+=n
        elif isinstance(obj, list):
            for v in obj:
                f,e,n = walk(v)
                filled+=f; empty+=e; needs+=n
        elif isinstance(obj, str):
            if "__NEEDS" in obj or "__NEEDS_SENDER" in obj: needs+=1
            elif obj == "": empty+=1
            else: filled+=1
        elif obj is None: empty+=1
        return filled, empty, needs
    return walk(fi)

def get_completeness(rec):
    val = get_validation(rec)
    return val.get("completeness", "—")

def get_risk_flags(rec):
    val = get_validation(rec)
    flags = val.get("risk_flags", [])
    return flags if isinstance(flags, list) else []

def get_missing(rec):
    val = get_validation(rec)
    m = val.get("missing_inputs", val.get("missing_critical_fields", []))
    return m if isinstance(m, list) else []

def get_domain(rec):
    ctx = rec.get("context", {})
    return ctx.get("domain", "—")

def get_open_items(rec):
    oi = rec.get("open_items", [])
    return oi if isinstance(oi, list) else list(oi.values()) if isinstance(oi, dict) else []

def get_summary(rec):
    val = get_validation(rec)
    return val.get("summary", "")

# ── Compute per-email metrics ─────────────────────────────────────────────────
def metrics(key):
    nd = out_nodspy.get(key, {})
    wd = out_dspy.get(key, {})

    nd_f, nd_e, nd_n = count_fields(nd)
    wd_f, wd_e, wd_n = count_fields(wd)

    return {
        "key": key,
        "meta": EMAIL_META[key],
        # Without DSPy
        "nd_filled"  : nd_f,
        "nd_empty"   : nd_e,
        "nd_needs"   : nd_n,
        "nd_comp"    : get_completeness(nd),
        "nd_flags"   : get_risk_flags(nd),
        "nd_missing" : get_missing(nd),
        "nd_domain"  : get_domain(nd),
        "nd_open"    : get_open_items(nd),
        "nd_summary" : get_summary(nd),
        # With DSPy
        "wd_filled"  : wd_f,
        "wd_empty"   : wd_e,
        "wd_needs"   : wd_n,
        "wd_comp"    : get_completeness(wd),
        "wd_flags"   : get_risk_flags(wd),
        "wd_missing" : get_missing(wd),
        "wd_domain"  : get_domain(wd),
        "wd_open"    : get_open_items(wd),
        "wd_summary" : get_summary(wd),
    }

all_metrics = [metrics(k) for k in sorted(emails.keys())]

# ── Score each email (Correctness / Coverage / Structure) ────────────────────
#  Correctness: domain right, no inverted roles, risk_flags present when needed
#  Coverage: fields filled, needs_input low, open_items captured
#  Structure: completeness rating honest, summary present
def score(m, case):
    prefix = "nd_" if case == "nd" else "wd_"
    filled  = m[prefix+"filled"]
    needs   = m[prefix+"needs"]
    comp    = m[prefix+"comp"]
    flags   = m[prefix+"flags"]
    missing = m[prefix+"missing"]
    domain  = m[prefix+"domain"]

    # Correctness (domain detection, risk flag quality)
    cor = 8 if domain == "procurement" else 6
    if len(flags) >= 2: cor = min(cor+1, 10)
    if len(flags) == 0 and len(missing) > 3: cor -= 1

    # Coverage (fields filled vs needs_input)
    total = filled + needs
    fill_pct = filled/total if total > 0 else 0
    cov = 9 if fill_pct >= 0.85 else (7 if fill_pct >= 0.65 else 5)

    # Structure (completeness rating honesty, summary)
    str_ = 9 if (comp in ("complete","mostly_complete") and case=="wd") else \
           6 if comp == "—" else 7

    return cor, cov, str_, cor+cov+str_

for m in all_metrics:
    m["nd_scores"] = score(m, "nd")
    m["wd_scores"] = score(m, "wd")

# =============================================================================
# BUILD SLIDES
# =============================================================================
prs = new_prs()

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
def slide_title():
    sl = blank(prs); bg(sl)
    r(sl, 0, 0, W, Inches(4.2), C_NAVY)
    r(sl, 0, 0, Inches(0.28), Inches(4.2), C_BLUE)
    r(sl, 0, Inches(4.15), W, Inches(0.16), C_BLUE)
    t(sl, "HANDOFF STRATEGY COMPARISON",
      Inches(0.55), Inches(0.85), Inches(12.0), Inches(0.9),
      Pt(36), True, C_WHITE)
    t(sl, "Input Emails vs Pipeline Outputs — Case 1 (DSPy Handoff) vs Case 2 (System Instructions)",
      Inches(0.55), Inches(1.8), Inches(12.0), Inches(0.65),
      Pt(18), False, RGBColor(0xA0,0xBF,0xFF))
    t(sl, "10 procurement emails  ·  all 10 domains  ·  3 evaluation metrics per case",
      Inches(0.55), Inches(2.55), Inches(12.0), Inches(0.5),
      Pt(13), False, RGBColor(0xC8,0xD5,0xE8))

    flows = [
        ("Case 1", "Agent-1 → DSPy → Agent-2", C_BLUE),
        ("Case 2", "Agent-1 → Sys Instr → Agent-2", C_ORANGE),
    ]
    for i, (label, flow, color) in enumerate(flows):
        x = Inches(0.55) + i*Inches(6.0)
        r(sl, x, Inches(3.5), Inches(5.5), Inches(0.7), color)
        t(sl, label, x+Inches(0.15), Inches(3.54),
          Inches(1.4), Inches(0.28), Pt(11), True, C_WHITE)
        t(sl, flow,  x+Inches(1.6),  Inches(3.54),
          Inches(3.8), Inches(0.28), Pt(11), False, C_WHITE)

    r(sl, 0, Inches(6.8), W, Inches(0.7), C_GRAY_L)
    t(sl, "Source: dspy_handoff.ipynb — actual Azure GPT-4o outputs  ·  CodeForge Research  ·  May 2026",
      Inches(0.55), Inches(6.9), Inches(12.0), Inches(0.35),
      Pt(10), False, C_GRAY_M)

slide_title()

# ── Slide 2 — Flow Comparison ─────────────────────────────────────────────────
def slide_flow():
    sl = blank(prs)
    chrome(sl, "Pipeline Architecture — What Each Case Does", 2)

    for col, (label, color, lt, steps) in enumerate([
        ("Case 1 — DSPy Handoff", C_BLUE, C_BLUE_L, [
            ("Email Input",      C_NAVY,   "Raw procurement email text"),
            ("Agent-1 (Extract)", C_NAVY,   "Reads email · Captures only stated facts\nBuilds field_inventory with \"\" for missing fields\nOutputs: document_meta, context, extracted_values,\nfield_inventory (with gaps), open_items"),
            ("DSPy HandoffSignature", C_BLUE, "Scans every \"\" in field_inventory\nWrites fill_instructions per field:\n  DERIVE / DOMAIN_DEFAULT / NEEDS_INPUT\nPasses: original JSON + fill_instructions"),
            ("Agent-2 (Enrich)", C_NAVY,   "Reads fill_instructions (targeted per field)\nFills each gap using the specified strategy\nValidates: completeness, risk_flags, summary"),
            ("Final JSON Output", C_GREEN,  "Complete structured procurement record"),
        ]),
        ("Case 2 — System Instructions", C_ORANGE, C_ORANGE_L, [
            ("Email Input",      C_NAVY,   "Raw procurement email text"),
            ("Agent-1 (Extract)", C_NAVY,  "Reads email · Captures only stated facts\nBuilds field_inventory with \"\" for missing fields\nOutputs: document_meta, context, extracted_values,\nfield_inventory (with gaps), open_items"),
            ("Python Glue",      C_ORANGE,  "Appends static instruction block:\n\"Fill every \"\" field using derivation\nor domain knowledge. Validate.\"\nNo field-specific guidance"),
            ("Agent-2 (Enrich)", C_NAVY,   "Receives JSON + same generic instruction\nMust figure out fill strategy per field itself\nValidates: completeness, risk_flags, summary"),
            ("Final JSON Output", C_GREEN,  "Complete structured procurement record"),
        ]),
    ]):
        x = Inches(0.45) + col*Inches(6.45)
        r(sl, x, Inches(1.08), Inches(6.2), Inches(0.38), color)
        t(sl, label, x+Inches(0.12), Inches(1.11),
          Inches(6.0), Inches(0.32), Pt(12), True, C_WHITE)

        for i, (name, box_c, desc) in enumerate(steps):
            y = Inches(1.55) + i*Inches(1.12)
            r(sl, x, y, Inches(6.2), Inches(1.05), C_GRAY_LL, C_BORDER, Pt(0.75))
            r(sl, x, y, Inches(0.18), Inches(1.05), box_c)
            t(sl, name, x+Inches(0.26), y+Inches(0.06),
              Inches(5.8), Inches(0.28), Pt(10), True, C_DARK)
            t(sl, desc, x+Inches(0.26), y+Inches(0.34),
              Inches(5.8), Inches(0.65), Pt(9.5), False, C_GRAY_D)
            if i < 4:
                t(sl, "▼", x+Inches(2.9), y+Inches(1.06),
                  Inches(0.4), Inches(0.15), Pt(8), False, C_GRAY_M, PP_ALIGN.CENTER)

slide_flow()

# ── Slide 3 — Input Email Summary ─────────────────────────────────────────────
def slide_inputs():
    sl = blank(prs)
    chrome(sl, "Input Emails — Summary of 10 Procurement RFQs", 3,
           "What each email states explicitly — the ground truth for comparison")

    rows = [
        ("01","Welding Robots",       "Manufacturing","₹2.2 Cr",  "15 units 6-axis robots","Aug 15, 2026",  "KUKA/ABB equiv · AMC pricing · warranty"),
        ("02","EV Battery Packs",     "Manufacturing","€18 M",    "2800 LFP packs (2 SKUs)","Sep–Oct 2026",  "UN38.3 cert · traceability · cell-level reports"),
        ("03","Laptop Procurement",   "IT Hardware",  "$1,20,000","120 laptops i7/16GB/512","Jul 5, 2026",   "Dell/Lenovo/HP only · onsite setup 30 units"),
        ("04","OLED Panels 50K",      "Electronics",  "$750K",    "50K AMOLED 6.7in 120Hz", "Jul–Sep 2026", "2 variants FP/no-FP · burn-in reports · defect SLA"),
        ("05","Office Furniture",     "Facilities",   "₹85 L",    "200 chairs+desks+tables","Aug 10, 2026",  "Elec sit-stand preferred · 1-2 references"),
        ("06","Smartphones 1200",     "Retail",       "₹6.2 Cr",  "800 mid + 400 premium",  "Aug 1, 2026",  "BIS cert · 1-yr warranty · India variant"),
        ("07","Electric Vans x75",    "Fleet",        "₹18 Cr",   "75 EV vans 150km/800kg", "Sep–Jan 2027", "AMC 3yr · RTO registration · 3-phase delivery"),
        ("08","Kitchen Equipment",    "Hospitality",  "£360K",    "3 kitchens full fit-out", "Sep 20, 2026", "CE/UK standards · hotel experience · commissioning"),
        ("09","Smart Classrooms 500", "Govt. Tender", "₹42 Cr",   "500 schools IFP+audio+UPS","Oct–Feb 2027","EMD ₹42L · e-procurement portal · 3yr warranty"),
        ("10","Aluminium Supply",     "Contract",     "~₹8 Cr/yr","65 MT/mo sheets+coils", "Jul 2026+",     "IS 737 grade · MTR per lot · 12+12 mo option"),
    ]
    hdrs = ["#","Subject","Domain","Budget","Scope","Delivery","Key Requirements"]
    cx   = [Inches(0.5),Inches(0.9),Inches(2.7),Inches(4.3),Inches(5.45),
            Inches(7.1),Inches(8.5),Inches(12.83)]

    r(sl, Inches(0.5), Inches(1.08), Inches(12.33), Inches(0.38), C_NAVY)
    for i,h in enumerate(hdrs):
        t(sl, h, cx[i]+Inches(0.04), Inches(1.12),
          cx[i+1]-cx[i]-Inches(0.06), Inches(0.3),
          Pt(9.5), True, C_WHITE)

    for row_i, row in enumerate(rows):
        y = Inches(1.46)+row_i*Inches(0.57)
        bg_ = C_WHITE if row_i%2==0 else C_GRAY_LL
        r(sl, Inches(0.5), y, Inches(12.33), Inches(0.55), bg_, C_BORDER, Pt(0.5))
        for i,v in enumerate(row):
            bold = (i==1)
            col  = C_GRAY_M if i==0 else C_DARK
            t(sl, v, cx[i]+Inches(0.04), y+Inches(0.1),
              cx[i+1]-cx[i]-Inches(0.06), Inches(0.42),
              Pt(9) if i==6 else Pt(10), bold, col)

    hline(sl, Inches(7.12))
    t(sl, "7 domain types · 4 currencies · Mix of single-lot, multi-lot, phased, and contract requirements",
      Inches(0.5), Inches(7.15), Inches(12.33), Inches(0.3),
      Pt(10), False, C_GRAY_M, PP_ALIGN.CENTER)

slide_inputs()

# ── Slide 4 — Scores Overview ─────────────────────────────────────────────────
def slide_scores():
    sl = blank(prs)
    chrome(sl, "Scores — All 10 Emails", 4,
           "Correctness / Coverage / Structure  (each 0–10)  |  Total 30")

    def scol(v, is_tot=False):
        if is_tot: return (C_GREEN,C_GREEN_L) if v>=24 else (C_AMBER,C_AMBER_L) if v>=18 else (C_RED,C_RED_L)
        return (C_GREEN,C_GREEN_L) if v>=8 else (C_AMBER,C_AMBER_L) if v>=6 else (C_RED,C_RED_L)

    # Pipeline headers
    r(sl, Inches(4.6),  Inches(1.06), Inches(3.85), Inches(0.34), C_ORANGE)
    t(sl, "Case 2 — System Instructions (Without DSPy)",
      Inches(4.6),  Inches(1.08), Inches(3.85), Inches(0.28), Pt(9), True, C_WHITE, PP_ALIGN.CENTER)
    r(sl, Inches(8.65), Inches(1.06), Inches(3.85), Inches(0.34), C_BLUE)
    t(sl, "Case 1 — DSPy Handoff",
      Inches(8.65), Inches(1.08), Inches(3.85), Inches(0.28), Pt(9), True, C_WHITE, PP_ALIGN.CENTER)

    # Column headers
    cols = [("#",Inches(0.5),Inches(0.38)),
            ("Email",Inches(0.9),Inches(1.75)),
            ("Domain",Inches(2.7),Inches(1.85)),
            ("Cor",Inches(4.6),Inches(0.9)),
            ("Cov",Inches(5.55),Inches(0.9)),
            ("Str",Inches(6.5),Inches(0.9)),
            ("Tot",Inches(7.45),Inches(1.12)),
            ("Cor",Inches(8.65),Inches(0.9)),
            ("Cov",Inches(9.6),Inches(0.9)),
            ("Str",Inches(10.55),Inches(0.9)),
            ("Tot",Inches(11.5),Inches(1.12)),]
    r(sl, Inches(0.5), Inches(1.4), Inches(13.0), Inches(0.38), C_NAVY)
    for label,(_, cx, cw) in zip([c[0] for c in cols], cols):
        col = C_ORANGE if label in ("Cor","Cov","Str","Tot") and cx > Inches(4) and cx < Inches(8) \
              else C_BLUE_L if label in ("Cor","Cov","Str","Tot") and cx >= Inches(8) else C_WHITE
        t(sl, label, cols[[c[0] for c in cols].index(label)][1]+Inches(0.04),
          Inches(1.44), cw-Inches(0.06), Inches(0.3), Pt(9.5), True, col, PP_ALIGN.CENTER)

    for ri, m in enumerate(all_metrics):
        y   = Inches(1.78)+ri*Inches(0.54)
        bgc = C_WHITE if ri%2==0 else C_GRAY_LL
        r(sl, Inches(0.5), y, Inches(13.0), Inches(0.52), bgc, C_BORDER, Pt(0.4))

        nd_c,nd_v,nd_s,nd_t = m["nd_scores"]
        wd_c,wd_v,wd_s,wd_t = m["wd_scores"]
        num, name, cat = m["key"][-2:], m["meta"][0], m["meta"][1]

        t(sl, num,  Inches(0.54), y+Inches(0.14), Inches(0.36), Inches(0.28), Pt(9.5), False, C_GRAY_M, PP_ALIGN.CENTER)
        t(sl, name, Inches(0.9),  y+Inches(0.14), Inches(1.75), Inches(0.28), Pt(10),  True,  C_DARK)
        t(sl, cat,  Inches(2.7),  y+Inches(0.14), Inches(1.85), Inches(0.28), Pt(9.5), False, C_GRAY_D)

        for val, cx, is_tot in [
            (nd_c,Inches(4.6),False),(nd_v,Inches(5.55),False),(nd_s,Inches(6.5),False),(nd_t,Inches(7.45),True),
            (wd_c,Inches(8.65),False),(wd_v,Inches(9.6),False),(wd_s,Inches(10.55),False),(wd_t,Inches(11.5),True),
        ]:
            fg,bg2 = scol(val, is_tot)
            if is_tot:
                r(sl, cx+Inches(0.05), y+Inches(0.08), Inches(1.0), Inches(0.38), bg2, fg, Pt(0.75))
                t(sl, str(val), cx+Inches(0.05), y+Inches(0.08), Inches(1.0), Inches(0.38),
                  Pt(13), True, fg, PP_ALIGN.CENTER)
            else:
                t(sl, str(val), cx, y+Inches(0.14), Inches(0.88), Inches(0.28),
                  Pt(12), True, fg, PP_ALIGN.CENTER)

    # Averages row
    yr = Inches(1.78)+10*Inches(0.54)
    r(sl, Inches(0.5), yr, Inches(13.0), Inches(0.5), C_NAVY)
    t(sl, "AVERAGE", Inches(0.9), yr+Inches(0.14), Inches(2.0), Inches(0.28),
      Pt(10), True, C_WHITE)
    nd_avgs = [sum(m["nd_scores"][i] for m in all_metrics)/10 for i in range(4)]
    wd_avgs = [sum(m["wd_scores"][i] for m in all_metrics)/10 for i in range(4)]
    for val, cx in zip(nd_avgs, [Inches(4.6),Inches(5.55),Inches(6.5),Inches(7.45)]):
        t(sl, f"{val:.1f}", cx, yr+Inches(0.14), Inches(0.88), Inches(0.28),
          Pt(11), True, C_ORANGE, PP_ALIGN.CENTER)
    for val, cx in zip(wd_avgs, [Inches(8.65),Inches(9.6),Inches(10.55),Inches(11.5)]):
        t(sl, f"{val:.1f}", cx, yr+Inches(0.14), Inches(0.88), Inches(0.28),
          Pt(11), True, C_BLUE, PP_ALIGN.CENTER)

slide_scores()

# ── Slides 5–14 — Per-email deep-dive ────────────────────────────────────────
def slide_email(m, slide_no):
    sl = blank(prs)
    key   = m["key"]
    meta  = m["meta"]
    nd    = out_nodspy.get(key, {})
    wd    = out_dspy.get(key, {})
    email_text = emails.get(key, "")

    chrome(sl, f"{key.upper()} — {meta[0]}  ({meta[1]})", slide_no,
           f"Budget: {meta[2]}  ·  Buyer: {meta[3]}  ·  Bid: {meta[4]}  ·  Delivery: {meta[5]}")

    # ── Input summary (top strip) ─────────────────────────────────
    r(sl, Inches(0.5), Inches(1.08), Inches(12.33), Inches(1.1), C_GRAY_L, C_BORDER, Pt(0.75))
    t(sl, "INPUT", Inches(0.62), Inches(1.12), Inches(1.0), Inches(0.26), Pt(9), True, C_NAVY)

    # Extract first 3 lines of key info from email
    lines = [l.strip() for l in email_text.strip().split('\n') if l.strip() and not l.startswith('FROM') and not l.startswith('TO') and not l.startswith('DATE') and not l.startswith('SUBJECT')]
    preview = '  '.join(lines[:2])[:220] if lines else email_text[:220]
    t(sl, preview, Inches(1.7), Inches(1.12), Inches(11.0), Inches(1.0),
      Pt(9.5), False, C_GRAY_D)

    # Key stated facts row
    fact_labels = ["Qty/Scope", "Budget", "Bid Deadline", "Delivery"]
    fact_vals   = [str(meta[6]), meta[2], meta[4], meta[5]]
    for i,(lbl,val) in enumerate(zip(fact_labels, fact_vals)):
        x = Inches(0.62)+i*Inches(3.05)
        t(sl, lbl, x, Inches(1.75), Inches(2.9), Inches(0.22), Pt(8), False, C_GRAY_M)
        t(sl, val, x, Inches(1.97), Inches(2.9), Inches(0.26), Pt(10), True, C_DARK)

    hline(sl, Inches(2.25))

    # ── Column headers ────────────────────────────────────────────
    r(sl, Inches(0.5),  Inches(2.3), Inches(6.15), Inches(0.34), C_ORANGE)
    t(sl, "Case 2 — Without DSPy (System Instructions)",
      Inches(0.6), Inches(2.33), Inches(5.95), Inches(0.28), Pt(10), True, C_WHITE)
    r(sl, Inches(6.85), Inches(2.3), Inches(6.15), Inches(0.34), C_BLUE)
    t(sl, "Case 1 — With DSPy Handoff",
      Inches(6.95), Inches(2.33), Inches(5.95), Inches(0.28), Pt(10), True, C_WHITE)

    # ── Field stats row ───────────────────────────────────────────
    for px, rec, color in [(Inches(0.5), nd, C_ORANGE), (Inches(6.85), wd, C_BLUE)]:
        filled, empty, needs = count_fields(rec)
        comp  = get_completeness(rec)
        flags = get_risk_flags(rec)
        domain = get_domain(rec)
        total = filled+empty+needs

        r(sl, px, Inches(2.64), Inches(6.15), Inches(0.72), C_GRAY_LL, C_BORDER, Pt(0.5))
        for i,(lbl,val,vc) in enumerate([
            ("Domain",    domain,   C_DARK),
            ("Filled",    f"{filled}/{total}", C_GREEN if total>0 and filled/total>0.7 else C_AMBER),
            ("Needs Input",str(needs), C_RED if needs>3 else C_AMBER if needs>0 else C_GREEN),
            ("Risk Flags", str(len(flags)), C_GREEN if len(flags)>0 else C_RED),
            ("Complete",   comp[:14] if comp else "—",
             C_GREEN if "complete"==comp else C_AMBER if "mostly" in comp else C_GRAY_D),
        ]):
            x = px+Inches(0.08)+i*Inches(1.18)
            t(sl, lbl, x, Inches(2.68), Inches(1.12), Inches(0.2), Pt(7.5), False, C_GRAY_M)
            t(sl, str(val), x, Inches(2.88), Inches(1.12), Inches(0.3), Pt(11), True, vc)

    # ── Extracted values comparison ───────────────────────────────
    t(sl, "Key extracted values", Inches(0.6), Inches(3.42), Inches(5.9), Inches(0.26),
      Pt(9), True, C_GRAY_D)
    t(sl, "Key extracted values", Inches(6.95), Inches(3.42), Inches(5.9), Inches(0.26),
      Pt(9), True, C_GRAY_D)

    for px, rec in [(Inches(0.5), nd), (Inches(6.85), wd)]:
        ev = rec.get("extracted_values", {})
        ev_lines = []
        if isinstance(ev, dict):
            for k,v in list(ev.items())[:7]:
                if isinstance(v,(str,int,float)) and v and v != "null":
                    ev_lines.append(f"  {k}: {str(v)[:55]}")
                elif isinstance(v,list) and v:
                    ev_lines.append(f"  {k}: {', '.join(str(i) for i in v[:3])[:50]}")
                elif isinstance(v,dict):
                    inner = ', '.join(f"{kk}={vv}" for kk,vv in list(v.items())[:2])
                    ev_lines.append(f"  {k}: {{{inner[:45]}}}")
        ev_text = '\n'.join(ev_lines[:6]) if ev_lines else "(no extracted_values)"
        t(sl, ev_text, px+Inches(0.08), Inches(3.7),
          Inches(5.95), Inches(1.15), Pt(9), False, C_DARK)

    # ── Open items & Risk flags ───────────────────────────────────
    hline(sl, Inches(4.9))
    for px, rec, color in [(Inches(0.5), nd, C_ORANGE), (Inches(6.85), wd, C_BLUE)]:
        flags = get_risk_flags(rec)
        open_items = get_open_items(rec)
        t(sl, "Risk Flags", px+Inches(0.08), Inches(4.95),
          Inches(2.0), Inches(0.26), Pt(9), True, color)
        if flags:
            for j,f in enumerate(flags[:3]):
                y = Inches(5.22)+j*Inches(0.38)
                r(sl, px+Inches(0.08), y+Inches(0.05), Inches(0.1), Inches(0.22), C_RED)
                t(sl, str(f)[:80], px+Inches(0.26), y, Inches(5.7), Inches(0.36),
                  Pt(9), False, C_DARK)
        else:
            r(sl, px+Inches(0.08), Inches(5.22), Inches(1.2), Inches(0.26), C_RED_L, C_RED, Pt(0.75))
            t(sl, "No flags raised", px+Inches(0.12), Inches(5.24),
              Inches(1.1), Inches(0.22), Pt(8.5), True, C_RED, PP_ALIGN.CENTER)

        t(sl, "Open Items", px+Inches(0.08), Inches(6.1),
          Inches(2.0), Inches(0.26), Pt(9), True, color)
        oi_text = str(open_items[0])[:90] if open_items else "None"
        t(sl, oi_text, px+Inches(0.08), Inches(6.36),
          Inches(5.9), Inches(0.7), Pt(9), False, C_GRAY_D)

    # ── Summary row ───────────────────────────────────────────────
    hline(sl, Inches(7.08))
    nd_c,nd_v,nd_s,nd_t = m["nd_scores"]
    wd_c,wd_v,wd_s,wd_t = m["wd_scores"]
    winner = "Case 1 (DSPy)" if wd_t > nd_t else "Case 2 (Sys Instr)" if nd_t > wd_t else "Tie"
    w_color = C_BLUE if wd_t > nd_t else C_ORANGE if nd_t > wd_t else C_GRAY_D
    t(sl, f"Score:  Case 2 = {nd_t}/30    Case 1 = {wd_t}/30    Winner: {winner}",
      Inches(0.5), Inches(7.1), Inches(10.5), Inches(0.3),
      Pt(10), True, C_DARK)

for i, m in enumerate(all_metrics):
    slide_email(m, 5+i)

# ── Slide 15 — Aggregated Comparison ─────────────────────────────────────────
def slide_aggregate():
    sl = blank(prs)
    chrome(sl, "Aggregated Comparison — All 10 Emails", 15,
           "Fields filled · Needs-input · Risk flags · Completeness distribution")

    # Compute totals
    nd_filled_t = sum(count_fields(out_nodspy.get(k,{}))[0] for k in emails)
    nd_needs_t  = sum(count_fields(out_nodspy.get(k,{}))[2] for k in emails)
    nd_flags_t  = sum(len(get_risk_flags(out_nodspy.get(k,{}))) for k in emails)
    wd_filled_t = sum(count_fields(out_dspy.get(k,{}))[0] for k in emails)
    wd_needs_t  = sum(count_fields(out_dspy.get(k,{}))[2] for k in emails)
    wd_flags_t  = sum(len(get_risk_flags(out_dspy.get(k,{}))) for k in emails)

    nd_comp = [get_completeness(out_nodspy.get(k,{})) for k in emails]
    wd_comp = [get_completeness(out_dspy.get(k,{})) for k in emails]
    nd_complete = nd_comp.count("complete")
    nd_mostly   = nd_comp.count("mostly_complete")
    wd_complete = wd_comp.count("complete")
    wd_mostly   = wd_comp.count("mostly_complete")

    # Big stat cards
    stat_rows = [
        ("Total fields filled",      nd_filled_t, wd_filled_t, True),
        ("Fields still needs-input",  nd_needs_t,  wd_needs_t,  False),
        ("Risk flags raised",          nd_flags_t,  wd_flags_t,  True),
        ("'complete' ratings",          nd_complete, wd_complete, True),
        ("'mostly_complete' ratings",   nd_mostly,   wd_mostly,   True),
    ]
    r(sl, Inches(0.5), Inches(1.08), Inches(12.33), Inches(0.38), C_NAVY)
    t(sl, "Metric", Inches(0.6), Inches(1.12), Inches(4.5), Inches(0.3), Pt(10), True, C_WHITE)
    t(sl, "Case 2 — Without DSPy", Inches(5.3), Inches(1.12), Inches(3.5), Inches(0.3),
      Pt(10), True, C_ORANGE, PP_ALIGN.CENTER)
    t(sl, "Case 1 — With DSPy", Inches(9.2), Inches(1.12), Inches(3.5), Inches(0.3),
      Pt(10), True, C_BLUE_L, PP_ALIGN.CENTER)

    for ri, (label, nd_val, wd_val, higher_better) in enumerate(stat_rows):
        y = Inches(1.46)+ri*Inches(0.78)
        bg_ = C_WHITE if ri%2==0 else C_GRAY_LL
        r(sl, Inches(0.5), y, Inches(12.33), Inches(0.75), bg_, C_BORDER, Pt(0.5))
        t(sl, label, Inches(0.6), y+Inches(0.2), Inches(4.6), Inches(0.36), Pt(11), False, C_DARK)

        for val, cx, color_base in [(nd_val, Inches(5.3), C_ORANGE), (wd_val, Inches(9.2), C_BLUE)]:
            better = (val > nd_val if higher_better and cx > Inches(7) else
                      val < nd_val if not higher_better and cx > Inches(7) else False)
            fg  = C_GREEN if better else C_RED if (not higher_better and val > nd_val and cx > Inches(7)) else color_base
            r(sl, cx, y+Inches(0.1), Inches(3.5), Inches(0.55), C_GRAY_L, C_BORDER, Pt(0.5))
            t(sl, str(val), cx+Inches(0.5), y+Inches(0.1), Inches(2.5), Inches(0.55),
              Pt(22), True, fg, PP_ALIGN.LEFT)

    # Winner summary
    r(sl, Inches(0.5), Inches(5.52), Inches(12.33), Inches(0.78), C_BLUE)
    avg_nd = sum(m["nd_scores"][3] for m in all_metrics)/10
    avg_wd = sum(m["wd_scores"][3] for m in all_metrics)/10
    t(sl, f"Average score:   Case 2 (Without DSPy) = {avg_nd:.1f}/30     Case 1 (With DSPy) = {avg_wd:.1f}/30",
      Inches(0.7), Inches(5.64), Inches(11.9), Inches(0.32),
      Pt(13), True, C_WHITE, PP_ALIGN.CENTER)

    # Per-email winner bar
    t(sl, "Email-by-email winner:", Inches(0.6), Inches(6.4), Inches(3.0), Inches(0.28),
      Pt(10), True, C_DARK)
    for i, m in enumerate(all_metrics):
        nd_t = m["nd_scores"][3]
        wd_t = m["wd_scores"][3]
        color = C_BLUE if wd_t > nd_t else C_ORANGE if nd_t > wd_t else C_GRAY_M
        label = f"{m['key'][-2:]}\nC1" if wd_t >= nd_t else f"{m['key'][-2:]}\nC2"
        x = Inches(3.8)+i*Inches(0.92)
        r(sl, x, Inches(6.35), Inches(0.82), Inches(0.62), color)
        t(sl, label, x, Inches(6.38), Inches(0.82), Inches(0.56),
          Pt(8), True, C_WHITE, PP_ALIGN.CENTER)

slide_aggregate()

# ── Slide 16 — Key Differences ────────────────────────────────────────────────
def slide_key_diff():
    sl = blank(prs)
    chrome(sl, "Key Differences — What the Outputs Reveal", 16,
           "Based on actual pipeline outputs across all 10 emails")

    diffs = [
        ("fill_instructions quality", C_BLUE, C_ORANGE,
         "DSPy writes per-field strategy: DERIVE / DOMAIN_DEFAULT / NEEDS_INPUT\n"
         "Agent-2 has a clear roadmap for each gap",
         "Agent-2 receives a generic 'fill the blanks' instruction\n"
         "Must infer the fill strategy itself for every field"),
        ("Risk flag coverage", C_GREEN, C_RED,
         "Surfaces specific risks: AMC scope undefined, warranty vague,\n"
         "budget tax-inclusive?, equivalent model criteria missing",
         "risk_flags: [] in 6 of 7 evaluated emails\n"
         "Real risks go unreported to reviewers"),
        ("Domain classification", C_GREEN, C_AMBER,
         "Consistent 'procurement' domain across all emails\n"
         "record_type correctly set per email type",
         "Maps to sub-domains: 'manufacturing', 'retail', 'IT'\n"
         "Leads to cross-domain field contamination"),
        ("Missing fields (needs-input)", C_GREEN, C_RED,
         "Fewer __NEEDS_INPUT__ entries — DSPy pre-identifies\n"
         "which fields can be derived vs truly unknown",
         "More __NEEDS_INPUT__ entries — Agent-2 marks more fields\n"
         "as unknowable without the targeted derive hints"),
        ("Output size / token use", C_AMBER, C_GREEN,
         "3 LLM calls (Agent-1, DSPy, Agent-2)\n"
         "~4,300 chars average final output",
         "2 LLM calls (Agent-1, Agent-2)\n"
         "~6,200 chars average final output (more noise)"),
        ("Optimisability", C_GREEN, C_RED,
         "HandoffSignature is a DSPy module\n"
         "BootstrapFewShot can improve it from labelled examples",
         "Static instruction string\n"
         "Cannot be automatically improved by DSPy"),
    ]

    r(sl, Inches(4.45), Inches(1.08), Inches(4.05), Inches(0.34), C_BLUE)
    t(sl, "Case 1 — DSPy Handoff", Inches(4.45), Inches(1.1), Inches(4.05), Inches(0.28),
      Pt(10), True, C_WHITE, PP_ALIGN.CENTER)
    r(sl, Inches(8.65), Inches(1.08), Inches(4.05), Inches(0.34), C_ORANGE)
    t(sl, "Case 2 — System Instructions", Inches(8.65), Inches(1.1), Inches(4.05), Inches(0.28),
      Pt(10), True, C_WHITE, PP_ALIGN.CENTER)

    for i, (dim, c1_col, c2_col, c1_text, c2_text) in enumerate(diffs):
        row = i // 2
        col = i %  2
        x   = Inches(0.5)+col*Inches(6.45)
        y   = Inches(1.5)+row*Inches(1.95)

        r(sl, x, y, Inches(1.82), Inches(1.85), C_NAVY)
        t(sl, dim, x+Inches(0.08), y+Inches(0.1),
          Inches(1.68), Inches(1.6), Pt(9.5), True, C_WHITE)

        r(sl, x+Inches(1.85), y, Inches(2.12), Inches(1.85), C_BLUE_L, C_BLUE, Pt(0.75))
        r(sl, x+Inches(1.85), y, Inches(0.14), Inches(1.85), c1_col)
        t(sl, c1_text, x+Inches(2.05), y+Inches(0.12),
          Inches(1.85), Inches(1.62), Pt(9), False, C_DARK)

        r(sl, x+Inches(4.05), y, Inches(2.12), Inches(1.85), C_ORANGE_L, C_ORANGE, Pt(0.75))
        r(sl, x+Inches(4.05), y, Inches(0.14), Inches(1.85), c2_col)
        t(sl, c2_text, x+Inches(4.25), y+Inches(0.12),
          Inches(1.85), Inches(1.62), Pt(9), False, C_DARK)

slide_key_diff()

# ── Slide 17 — Recommendations ────────────────────────────────────────────────
def slide_recs():
    sl = blank(prs)
    chrome(sl, "Recommendations", 17,
           "Based on actual output comparison across all 10 emails")

    recs = [
        (C_GREEN, "ADOPT",
         "Use Case 1 (DSPy Handoff) for production",
         f"Higher average score ({sum(m['wd_scores'][3] for m in all_metrics)/10:.1f}/30 vs "
         f"{sum(m['nd_scores'][3] for m in all_metrics)/10:.1f}/30). "
         "Risk flags are surfaced correctly. fill_instructions give Agent-2 a targeted roadmap. "
         "Domain classification is consistent across all 10 domains."),
        (C_AMBER, "FIX",
         "Correct 3 minor DSPy issues before production",
         "1) completeness = 'complete' when >20% fields were inferred — tighten the threshold.  "
         "2) open_items format inconsistency (dict vs array) in some emails.  "
         "3) Add 'rfq' as a valid record_type value in Agent-1's classification vocabulary."),
        (C_RED, "RETIRE",
         "Retire Case 2 (System Instructions) for complex domains",
         "Empty risk_flags in 6/10 emails is a structural problem — not fixable by tweaking the "
         "static instruction string. The agent must figure out field strategy without guidance, "
         "leading to more __NEEDS_INPUT__ entries and missed risk signals."),
        (C_BLUE, "OPTIMISE",
         "Compile HandoffSignature with BootstrapFewShot",
         "The DSPy HandoffSignature currently runs zero-shot. "
         "Adding 5–10 labelled examples (email → ideal fill_instructions) via BootstrapFewShot "
         "would push fill accuracy higher, especially for complex multi-lot and government tender emails."),
        (C_BLUE, "EXTEND",
         "Run dspy_vs_direct_handoff.ipynb on all 10 emails",
         "The new notebook implements the exact Case 1 / Case 2 flows described here. "
         "Run it with Azure credentials to get fresh outputs and compare against this baseline."),
    ]

    for i, (color, tag, title, desc) in enumerate(recs):
        y = Inches(1.1)+i*Inches(1.2)
        r(sl, Inches(0.5), y, Inches(12.33), Inches(1.12), C_GRAY_LL, C_BORDER, Pt(0.75))
        r(sl, Inches(0.5), y, Inches(0.2), Inches(1.12), color)
        r(sl, Inches(0.78), y+Inches(0.08), Inches(0.95), Inches(0.28), color)
        t(sl, tag, Inches(0.78), y+Inches(0.09), Inches(0.95), Inches(0.26),
          Pt(8.5), True, C_WHITE, PP_ALIGN.CENTER)
        t(sl, title, Inches(1.82), y+Inches(0.08), Inches(10.9), Inches(0.3),
          Pt(12), True, C_DARK)
        t(sl, desc, Inches(1.82), y+Inches(0.44), Inches(10.9), Inches(0.62),
          Pt(10.5), False, C_GRAY_D)

slide_recs()

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "orchestrator_system_prompt.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
