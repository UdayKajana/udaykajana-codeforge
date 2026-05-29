"""
3-slide Prodapt-branded PPT: DSPy vs 2-Agent Direct on real Japanese email.
Uses pod_copy- Copy.pptx as the template (inherits Prodapt master background).
All metrics from actual notebook run on japanese_email_test.json (Suzuki Sekkei).
Run:  python generate_prodapt_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

# ── Prodapt brand palette ────────────────────────────────────────────────────
P_RED      = RGBColor(0xEB, 0x26, 0x2A)   # Prodapt primary red
P_DARK     = RGBColor(0x1C, 0x1C, 0x1C)   # near-black
P_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
P_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)   # card backgrounds
P_MGRAY    = RGBColor(0xAB, 0xAB, 0xAB)   # labels / secondary text
P_DGRAY    = RGBColor(0x4A, 0x4A, 0x4A)   # body text
P_GREEN    = RGBColor(0x1B, 0x87, 0x3E)   # correct / positive
P_GREEN_LT = RGBColor(0xE6, 0xF7, 0xEC)
P_RED_LT   = RGBColor(0xFD, 0xE8, 0xE8)
P_AMBER    = RGBColor(0xB3, 0x5A, 0x00)
P_AMBER_LT = RGBColor(0xFE, 0xF3, 0xE0)
P_BORDER   = RGBColor(0xDC, 0xDC, 0xDC)
P_BLUE_LT  = RGBColor(0xE8, 0xF0, 0xFE)
P_BLUE     = RGBColor(0x18, 0x65, 0xD6)

W, H = Inches(13.33), Inches(7.5)

# ── Open template ────────────────────────────────────────────────────────────
prs = Presentation("pod_copy- Copy.pptx")

# Remove the placeholder slide that ships with the template
xml_slides = prs.slides._sldIdLst
while len(xml_slides):
    rId = xml_slides[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[0])

BLANK = prs.slide_layouts[42]   # 'Blank slide'

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=P_WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=Pt(11), bold=False, color=P_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = size
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def header_bar(slide, title, subtitle=""):
    """Red Prodapt header band at top of slide."""
    rect(slide, 0, 0, W, Inches(1.0), P_RED)
    txt(slide, title,
        Inches(0.45), Inches(0.12), Inches(11.5), Inches(0.52),
        Pt(20), True, P_WHITE, PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.45), Inches(0.64), Inches(11.5), Inches(0.3),
            Pt(10), False, RGBColor(0xFF, 0xC0, 0xBF), PP_ALIGN.LEFT)

def metric_pill(slide, x, y, label, value, good=None):
    """Small metric card: label on top, value big below."""
    bg = P_RED_LT if good is False else (P_GREEN_LT if good is True else P_LGRAY)
    fg = P_RED    if good is False else (P_GREEN    if good is True else P_DARK)
    bord = P_RED  if good is False else (P_GREEN    if good is True else P_BORDER)
    rect(slide, x, y, Inches(2.8), Inches(1.05), bg, bord, Pt(1.0))
    txt(slide, label, x + Inches(0.14), y + Inches(0.06),
        Inches(2.55), Inches(0.28), Pt(9), False, P_MGRAY)
    txt(slide, value, x + Inches(0.14), y + Inches(0.36),
        Inches(2.55), Inches(0.55), Pt(26), True, fg)


# =============================================================================
# SLIDE 1 — The Problem  (Pipeline B: 2-Agent Direct, genuine results)
# =============================================================================

def slide_problem():
    sl = add_slide()
    fill_bg(sl)
    header_bar(sl,
               "The Problem: 2-Agent Direct on a Real Japanese Procurement Email",
               "Schema inflation, missing domain fills, and wrong risk flags — genuine results from notebook run")

    EMAIL_W = Inches(5.4)
    EMAIL_X = Inches(0.42)
    CONTENT_Y = Inches(1.12)

    # ── Left: actual test email (Suzuki Sekkei) ──────────────────────────────
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(4.65), P_LGRAY, P_BORDER, Pt(0.75))
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(0.34), P_DARK)
    txt(sl, "ACTUAL TEST EMAIL  (fully Japanese — Suzuki Sekkei RFQ)",
        EMAIL_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        EMAIL_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    email_lines = [
        ("From:",    "yamamoto.kenta@suzuki-sekkei.co.jp"),
        ("To:",      "bids@globalparts.com"),
        ("Subject:", "[Urgent] 6-axis arc welding robot RFQ"),
        ("Date:",    "2026-05-26  (Monday 10:32 JST)"),
        ("",         ""),
        ("",         "Suzuki Sekkei Co., Ltd. [鈴木設計株式会社]"),
        ("",         "18 units 6-axis arc welding robots"),
        ("",         "FANUC M-10iA/12 or Yaskawa Motoman MA1440"),
        ("Budget:",  "JPY 220,000,000 (excl. tax) — delivery incl."),
        ("",         ""),
        ("Phase 1:", "6 units  Sep 15, 2026"),
        ("Phase 2:", "6 units  Nov 30, 2026"),
        ("Phase 3:", "6 units  Jan 31, 2027"),
        ("",         ""),
        ("AMC:",     "Min. 3 years — scope UNDEFINED in email"),
        ("Bid by:",  "June 20, 2026  17:00 JST"),
        ("Contact:", "Yamamoto Kenta [山本健太]  — Procurement Mgr"),
    ]
    for i, (k, v) in enumerate(email_lines):
        y = CONTENT_Y + Inches(0.44) + i * Inches(0.245)
        if k:
            txt(sl, k, EMAIL_X + Inches(0.14), y,
                Inches(0.82), Inches(0.22), Pt(8.5), True, P_DGRAY)
            txt(sl, v, EMAIL_X + Inches(0.98), y,
                EMAIL_W - Inches(1.1), Inches(0.22), Pt(8.5), False, P_DARK)
        else:
            txt(sl, v, EMAIL_X + Inches(0.14), y,
                EMAIL_W - Inches(0.3), Inches(0.22), Pt(8.5), False, P_DARK)

    # ── Right: 2-Agent Direct output (actual errors from notebook) ────────────
    OUT_X = Inches(6.08)
    OUT_W = Inches(6.82)
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(4.65), P_WHITE, P_BORDER, Pt(0.75))
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(0.34), P_AMBER)
    txt(sl, "2-AGENT DIRECT OUTPUT  (Pipeline B — actual notebook output)",
        OUT_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        OUT_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    # Actual errors from eval_results.py analysis of Pipeline B
    errors = [
        ("responding_party",  "6x NEEDS_INPUT",             "SCHEMA INFLATION", P_RED,
         "Vendor data has no place in a buyer-side RFQ — phantom section created"),
        ("secondary_contact", "6x NEEDS_INPUT",             "SCHEMA INFLATION",  P_RED,
         "12 spurious NEEDS_INPUT total — only 2 are genuine gaps in the email"),
        ("risk_flags",        '"No supplier details..."',   "WRONG RISK FLAG",   P_RED,
         "Supplier details are never in a buyer-issued RFQ by definition"),
        ("payment_terms",     "null  (left blank)",          "MISSING DEFAULT",   P_AMBER,
         "Net 30 (industry standard) not applied — ERP record incomplete"),
        ("governing_law",     "null  (left blank)",          "MISSING DEFAULT",   P_AMBER,
         "Laws of Japan not added — 5 more domain defaults also missing"),
        ("fields_filled",     "45  (of 55 expected)",        "LOW COVERAGE",      P_AMBER,
         "No domain defaults applied — record needs manual completion before ERP use"),
    ]

    for i, (field, value, badge, bc, note) in enumerate(errors):
        y = CONTENT_Y + Inches(0.46) + i * Inches(0.67)
        bg = P_RED_LT if bc == P_RED else P_AMBER_LT
        rect(sl, OUT_X + Inches(0.12), y, OUT_W - Inches(0.22), Inches(0.62), bg)
        txt(sl, field + ":", OUT_X + Inches(0.22), y + Inches(0.06),
            Inches(1.7), Inches(0.26), Pt(9), True, P_DARK)
        txt(sl, value, OUT_X + Inches(1.95), y + Inches(0.06),
            Inches(2.1), Inches(0.26), Pt(9), False, P_DARK, PP_ALIGN.LEFT, True)
        bw = Inches(0.14) + Pt(7) * len(badge) / Pt(1) * Inches(0.0125)
        rect(sl, OUT_X + Inches(4.15), y + Inches(0.06), bw, Inches(0.24), bc)
        txt(sl, badge, OUT_X + Inches(4.18), y + Inches(0.07),
            bw - Inches(0.06), Inches(0.2), Pt(7.5), True, P_WHITE, PP_ALIGN.LEFT)
        txt(sl, note, OUT_X + Inches(0.22), y + Inches(0.34),
            OUT_W - Inches(0.5), Inches(0.22), Pt(8.5), False, P_DGRAY,
            PP_ALIGN.LEFT, True, True)

    # ── Bottom: Pipeline B real metric pills ─────────────────────────────────
    rect(sl, Inches(0.42), Inches(5.88), W - Inches(0.84), Inches(0.28), P_DARK)
    txt(sl, "PIPELINE B METRICS  (2-Agent Direct — Suzuki Sekkei Japanese RFQ)",
        Inches(0.55), Inches(5.9), Inches(8.0), Inches(0.24),
        Pt(9.5), True, P_WHITE)

    # Genuine Pipeline B scores from notebook run
    metrics = [
        ("F1 Score  (extraction accuracy)", "0.83", False),
        ("Hallucination Rate",              "15%",  False),
        ("Confidence Score",                "64%",  False),
        ("Faithfulness",                    "0.95", True),   # B wins this one
    ]
    gap = (W - Inches(0.84) - 4 * Inches(2.8)) / 3
    for i, (lbl, val, good) in enumerate(metrics):
        x = Inches(0.42) + i * (Inches(2.8) + gap)
        metric_pill(sl, x, Inches(6.24), lbl, val, good)


# =============================================================================
# SLIDE 2 — DSPy in Action  (Pipeline A: DSPy Orchestrator, genuine results)
# =============================================================================

def slide_dspy():
    sl = add_slide()
    fill_bg(sl)
    header_bar(sl,
               "DSPy Orchestrator in Action: Same Email, Structured Audit",
               "Genuine results: 55 fields filled, 7 domain defaults, only 2 real NEEDS_INPUT")

    EMAIL_W = Inches(5.4)
    EMAIL_X = Inches(0.42)
    CONTENT_Y = Inches(1.12)

    # ── Left: same email (key fields) ────────────────────────────────────────
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(2.88), P_LGRAY, P_BORDER, Pt(0.75))
    rect(sl, EMAIL_X, CONTENT_Y, EMAIL_W, Inches(0.34), P_DARK)
    txt(sl, "SAME EMAIL  (Suzuki Sekkei — japanese_email_test.json)",
        EMAIL_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        EMAIL_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    email_lines = [
        ("From:",    "yamamoto.kenta@suzuki-sekkei.co.jp"),
        ("Subject:", "[Urgent] 6-axis arc welding robot RFQ"),
        ("",         "18 units — FANUC M-10iA/12 or Yaskawa MA1440"),
        ("Budget:",  "JPY 220,000,000 (excl. tax) — delivery incl."),
        ("Phases:",  "6+6+6 units: Sep/Nov 2026, Jan 2027"),
        ("AMC:",     "Min. 3 years — scope NOT defined in email"),
        ("Bid by:",  "June 20, 2026  17:00 JST"),
        ("Contact:", "Yamamoto Kenta [山本健太] — Procurement Mgr"),
    ]
    for i, (k, v) in enumerate(email_lines):
        y = CONTENT_Y + Inches(0.44) + i * Inches(0.295)
        if k:
            txt(sl, k, EMAIL_X + Inches(0.14), y,
                Inches(0.82), Inches(0.26), Pt(8.5), True, P_DGRAY)
            txt(sl, v, EMAIL_X + Inches(0.98), y,
                EMAIL_W - Inches(1.1), Inches(0.26), Pt(8.5), False, P_DARK)
        else:
            txt(sl, v, EMAIL_X + Inches(0.14), y,
                EMAIL_W - Inches(0.3), Inches(0.26), Pt(8.5), False, P_DARK)

    # Risk flags box — what DSPy found (but note: risk_flags was empty, critical gap)
    rect(sl, EMAIL_X, Inches(4.1), EMAIL_W, Inches(1.14), P_AMBER_LT, P_AMBER, Pt(0.75))
    txt(sl, "Real Risks in Email  (NEITHER pipeline fully flagged all 3)",
        EMAIL_X + Inches(0.14), Inches(4.15), EMAIL_W - Inches(0.25), Inches(0.26),
        Pt(9), True, P_AMBER)
    flags = [
        "AMC scope and inclusions completely undefined in email",
        "Budget excludes 2% consumption tax — clarification needed in bid",
        "No RFQ reference number — bid tracking at risk",
    ]
    for i, f in enumerate(flags):
        txt(sl, f"X  {f}", EMAIL_X + Inches(0.14), Inches(4.46) + i * Inches(0.24),
            EMAIL_W - Inches(0.28), Inches(0.22), Pt(8.5), False, P_DARK)

    # ── Right: DSPy Pipeline A output (correct fields) ───────────────────────
    OUT_X = Inches(6.08)
    OUT_W = Inches(6.82)
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(3.94), P_WHITE, P_BORDER, Pt(0.75))
    rect(sl, OUT_X, CONTENT_Y, OUT_W, Inches(0.34), P_GREEN)
    txt(sl, "DSPY ORCHESTRATOR OUTPUT  (Pipeline A — actual extracted values)",
        OUT_X + Inches(0.12), CONTENT_Y + Inches(0.05),
        OUT_W - Inches(0.2), Inches(0.28), Pt(9.5), True, P_WHITE)

    # Actual Pipeline A outputs from notebook
    correct = [
        ("buyer_company",   "Suzuki Sekkei Co., Ltd.  [鈴木設計株式会社]"),
        ("budget",          "JPY 220,000,000 (excl. tax) — delivery incl., AMC separate"),
        ("quantity",        "18 units — 3-phase delivery (6+6+6)"),
        ("delivery_phase1", "September 15, 2026  [2026年9月15日]"),
        ("payment_terms",   "Net 30  (industry standard — domain default applied)"),
        ("governing_law",   "Laws of Japan  (domain default applied)"),
        ("contact_email",   "yamamoto.kenta@suzuki-sekkei.co.jp"),
        ("robot_model",     "FANUC M-10iA/12 or Yaskawa Motoman MA1440"),
        ("bid_deadline",    "June 20, 2026  17:00 JST  [2026年6月20日]"),
    ]

    for i, (field, value) in enumerate(correct):
        y = CONTENT_Y + Inches(0.46) + i * Inches(0.375)
        bg = P_WHITE if i % 2 == 0 else P_GREEN_LT
        rect(sl, OUT_X + Inches(0.12), y, OUT_W - Inches(0.22), Inches(0.36), bg)
        txt(sl, field + ":", OUT_X + Inches(0.22), y + Inches(0.05),
            Inches(1.65), Inches(0.26), Pt(8.5), True, P_DARK)
        txt(sl, value, OUT_X + Inches(1.92), y + Inches(0.05),
            OUT_W - Inches(2.4), Inches(0.26), Pt(8.5), False, P_DARK)
        txt(sl, "OK", OUT_X + OUT_W - Inches(0.36), y + Inches(0.05),
            Inches(0.3), Inches(0.26), Pt(8), True, P_GREEN, PP_ALIGN.CENTER)

    # ── Comparison table (bottom right) — genuine numbers ────────────────────
    T_X, T_Y = OUT_X, Inches(5.18)
    T_W, T_H = OUT_W, Inches(1.96)

    col_x = [T_X, T_X + Inches(2.22), T_X + Inches(4.22), T_X + Inches(6.04)]
    col_w = [Inches(2.18), Inches(1.96), Inches(1.78), Inches(0.74)]

    rect(sl, T_X, T_Y, T_W, Inches(0.34), P_DARK)
    for label, cx, cw in zip(
        ["Metric", "2-Agent Direct", "DSPy Orchestrator", "Delta"], col_x, col_w
    ):
        align = PP_ALIGN.LEFT if label == "Metric" else PP_ALIGN.CENTER
        txt(sl, label, cx + Inches(0.06), T_Y + Inches(0.06),
            cw - Inches(0.08), Inches(0.26), Pt(9), True, P_WHITE, align)

    # Actual comparison: Pipeline B vs Pipeline A
    table_rows = [
        ("F1 Score",         "0.83", "0.87", "+5%",   True),
        ("Hallucination",    "15%",  "8%",   "-47%",  True),
        ("Confidence Score", "64%",  "79%",  "+23%",  True),
        ("Faithfulness",     "0.95", "0.92", "-3%*",  False),  # B wins faithfulness
    ]
    for r, (metric, before, after, delta, dspy_wins) in enumerate(table_rows):
        y = T_Y + Inches(0.34) + r * Inches(0.4)
        bg = P_WHITE if r % 2 == 0 else P_LGRAY
        rect(sl, T_X, y, T_W, Inches(0.38), bg, P_BORDER, Pt(0.4))
        txt(sl, metric, col_x[0] + Inches(0.06), y + Inches(0.08),
            col_w[0] - Inches(0.08), Inches(0.26), Pt(9.5), True, P_DARK)
        txt(sl, before, col_x[1] + Inches(0.06), y + Inches(0.08),
            col_w[1] - Inches(0.08), Inches(0.26), Pt(10), True, P_RED, PP_ALIGN.CENTER)
        after_color = P_GREEN if dspy_wins else P_AMBER
        txt(sl, after, col_x[2] + Inches(0.06), y + Inches(0.08),
            col_w[2] - Inches(0.08), Inches(0.26), Pt(10), True, after_color, PP_ALIGN.CENTER)
        delta_bg = P_GREEN if dspy_wins else P_AMBER
        rect(sl, col_x[3], y + Inches(0.06), col_w[3], Inches(0.26), delta_bg)
        txt(sl, delta, col_x[3], y + Inches(0.06),
            col_w[3], Inches(0.26), Pt(9), True, P_WHITE, PP_ALIGN.CENTER)

    # footnote for faithfulness
    txt(sl, "* Faithfulness 0.95 vs 0.92: 2-Agent Direct wins (readable dates, Japanese traceability brackets)",
        T_X + Inches(0.06), T_Y + Inches(1.62),
        T_W - Inches(0.1), Inches(0.28), Pt(7.5), False, P_MGRAY, PP_ALIGN.LEFT, True)

    # Left bottom — confidence summary
    rect(sl, EMAIL_X, Inches(5.18), EMAIL_W, Inches(0.56), P_GREEN_LT, P_GREEN, Pt(0.75))
    txt(sl, "Pipeline A (DSPy):",
        EMAIL_X + Inches(0.14), Inches(5.24),
        Inches(2.0), Inches(0.38), Pt(9), True, P_GREEN)
    txt(sl, "79% Confidence  |  55 Fields  |  7 Domain Defaults  |  2 Genuine NEEDS_INPUT",
        EMAIL_X + Inches(2.2), Inches(5.24),
        EMAIL_W - Inches(2.35), Inches(0.38), Pt(9), True, P_GREEN)


# =============================================================================
# SLIDE 3 — Summary  (genuine findings)
# =============================================================================

def slide_summary():
    sl = add_slide()
    fill_bg(sl)
    header_bar(sl,
               "Summary",
               "Genuine findings from Japanese RFQ test — Suzuki Sekkei Co., Ltd. [鈴木設計株式会社]")

    MID = Inches(6.67)   # centre divider
    COL_Y = Inches(1.16)

    # ── Column headers ────────────────────────────────────────────────────────
    rect(sl, Inches(0.42), COL_Y, Inches(5.92), Inches(0.38), P_RED)
    txt(sl, "Current Challenges  (2-Agent Direct — Pipeline B)",
        Inches(0.55), COL_Y + Inches(0.07),
        Inches(5.7), Inches(0.28), Pt(11), True, P_WHITE)

    rect(sl, MID + Inches(0.15), COL_Y, Inches(6.08), Inches(0.38), P_GREEN)
    txt(sl, "How DSPy Addresses Them  (Pipeline A)",
        MID + Inches(0.28), COL_Y + Inches(0.07),
        Inches(5.8), Inches(0.28), Pt(11), True, P_WHITE)

    # ── Divider ───────────────────────────────────────────────────────────────
    rect(sl, MID + Inches(0.06), COL_Y, Inches(0.06), Inches(5.38), P_BORDER)

    # Genuine findings from eval_results.py and compare3.py
    challenges = [
        ("Schema Inflation: 12 Phantom Fields",
         "2-Agent Direct created responding_party + secondary_contact sections "
         "with 12 NEEDS_INPUT for vendor data that does not belong in a buyer-side RFQ. "
         "Only 2 of 14 reported NEEDS_INPUT were real gaps. Confidence: 64%."),
        ("No Domain Intelligence Applied",
         "Zero domain defaults added — payment_terms, governing_law, JCAA arbitration, "
         "force_majeure, bid_validity_days all left blank. Only 45 of 55 expected fields "
         "filled. Record needs manual completion before ERP processing."),
        ("Risk Blind Spot — Critical for Both Pipelines",
         "Neither pipeline flagged all 3 real procurement risks: undefined AMC scope, "
         "budget excludes 2% consumption tax, no RFQ reference number. "
         "This is the #1 improvement target before any production deployment."),
    ]

    solutions = [
        ("DSPy Audit Eliminates Phantom Fields",
         "AuditSignature validates every field against the source email. "
         "Zero phantom sections created — only 2 genuine NEEDS_INPUT (rfq_id, ref no.) "
         "vs 12 self-inflicted by 2-Agent Direct. "
         "F1 Score: 0.87 vs 0.83. Confidence: 79% vs 64%."),
        ("7 Domain Defaults Auto-Applied",
         "DSPy fills industry-standard fields automatically: Net 30 payment terms, "
         "Laws of Japan, JCAA arbitration, force_majeure clause, bid_validity_days (90 days). "
         "55 fields filled vs 45 — ERP record is actionable today. +15pp confidence gain."),
        ("Next Step: Dedicated Risk Signature",
         "DSPy reduces hallucination to 8% (vs 15%) and wins F1, Confidence, and "
         "Hallucination metrics. Faithfulness: 2-Agent wins (0.95 vs 0.92) via readable "
         "dates and Japanese brackets. Add RiskSignature to catch AMC, tax, and ref gaps."),
    ]

    for i, ((c_title, c_body), (s_title, s_body)) in enumerate(zip(challenges, solutions)):
        y = COL_Y + Inches(0.5) + i * Inches(1.6)
        row_h = Inches(1.48)

        # Challenge card (left, red)
        rect(sl, Inches(0.42), y, Inches(5.92), row_h, P_RED_LT, P_RED, Pt(0.75))
        rect(sl, Inches(0.42), y, Inches(0.2), row_h, P_RED)
        txt(sl, c_title, Inches(0.72), y + Inches(0.1),
            Inches(5.4), Inches(0.3), Pt(12), True, P_DARK)
        txt(sl, c_body, Inches(0.72), y + Inches(0.46),
            Inches(5.4), Inches(0.9), Pt(10), False, P_DGRAY)

        # Solution card (right, green)
        rect(sl, MID + Inches(0.15), y, Inches(6.08), row_h, P_GREEN_LT, P_GREEN, Pt(0.75))
        rect(sl, MID + Inches(0.15), y, Inches(0.2), row_h, P_GREEN)
        txt(sl, s_title, MID + Inches(0.45), y + Inches(0.1),
            Inches(5.6), Inches(0.3), Pt(12), True, P_DARK)
        txt(sl, s_body, MID + Inches(0.45), y + Inches(0.46),
            Inches(5.6), Inches(0.9), Pt(10), False, P_DGRAY)

    # ── Bottom KPI bar — genuine numbers ──────────────────────────────────────
    rect(sl, Inches(0.42), Inches(6.04), W - Inches(0.84), Inches(0.38), P_RED)
    txt(sl, "GENUINE RESULTS  (Japanese RFQ — Suzuki Sekkei)",
        Inches(0.55), Inches(6.08), Inches(5.0), Inches(0.28),
        Pt(10), True, P_WHITE)

    # Real numbers: B -> A
    kpis = [
        ("F1 Score",     "0.83  to  0.87",     "+5%",    P_GREEN),
        ("Hallucination","15%  to  8%",         "-47%",   P_GREEN),
        ("Confidence",   "64%  to  79%",        "+23%",   P_GREEN),
        ("Faithfulness", "0.95 (B) / 0.92 (A)", "B +3pp", P_AMBER),  # B wins
        ("Integration",  "OpenAI Agent SDK",    "Ready",  P_BLUE),
    ]
    seg_w = (W - Inches(0.84)) / len(kpis)
    for i, (lbl, val, delta, color) in enumerate(kpis):
        x = Inches(0.42) + i * seg_w
        rect(sl, x, Inches(6.46), seg_w - Inches(0.06), Inches(0.9),
             P_WHITE, P_BORDER, Pt(0.5))
        txt(sl, lbl, x + Inches(0.1), Inches(6.5),
            seg_w - Inches(0.18), Inches(0.24), Pt(8.5), False, P_MGRAY)
        txt(sl, val, x + Inches(0.1), Inches(6.74),
            seg_w - Inches(0.18), Inches(0.3), Pt(9.5), True, P_DARK)
        rect(sl, x + seg_w - Inches(0.64), Inches(6.5),
             Inches(0.58), Inches(0.24), color)
        txt(sl, delta, x + seg_w - Inches(0.64), Inches(6.51),
            Inches(0.58), Inches(0.22), Pt(7.5), True, P_WHITE, PP_ALIGN.CENTER)


# ── Build ────────────────────────────────────────────────────────────────────
slide_problem()
slide_dspy()
slide_summary()

OUT = "prodapt_dspy_ppt.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
