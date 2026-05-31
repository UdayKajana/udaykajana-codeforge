"""
DSPy Evaluation v2 PPT — 4 slides
  Slide 1: Cover (same)
  Slide 2: Email (left) + Agent flow diagram + DSPy intro (right)
  Slide 3: No-DSPy metrics table (per-metric: score, went-well, did-not)
  Slide 4: DSPy metrics table + Final summary comparison table
Run:  python generate_ver2_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Prodapt brand palette ────────────────────────────────────────────────────
P_RED      = RGBColor(0xEB, 0x26, 0x2A)
P_DARK     = RGBColor(0x1C, 0x1C, 0x1C)
P_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
P_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
P_MGRAY    = RGBColor(0xAB, 0xAB, 0xAB)
P_DGRAY    = RGBColor(0x4A, 0x4A, 0x4A)
P_GREEN    = RGBColor(0x1B, 0x87, 0x3E)
P_GREEN_LT = RGBColor(0xE6, 0xF7, 0xEC)
P_RED_LT   = RGBColor(0xFD, 0xE8, 0xE8)
P_AMBER    = RGBColor(0xB3, 0x5A, 0x00)
P_AMBER_LT = RGBColor(0xFE, 0xF3, 0xE0)
P_BORDER   = RGBColor(0xDC, 0xDC, 0xDC)
P_BLUE     = RGBColor(0x18, 0x65, 0xD6)
P_PURPLE   = RGBColor(0x6B, 0x2F, 0x9E)
P_PURP_LT  = RGBColor(0xF0, 0xE8, 0xFA)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation("pod_copy- Copy.pptx")
xml_slides = prs.slides._sldIdLst
while len(xml_slides):
    rId = xml_slides[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[0])

COVER = prs.slide_layouts[1]   # Cover with two logos
BLANK = prs.slide_layouts[42]  # Blank slide

# ── Core helpers ─────────────────────────────────────────────────────────────

def bg(slide, color=P_WHITE):
    b = slide.background
    b.fill.solid()
    b.fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h,
       sz=Pt(10), bold=False, color=P_DARK,
       align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = sz
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return t

def hdr(slide, title, subtitle=""):
    box(slide, 0, 0, W, Inches(1.0), P_RED)
    tb(slide, title,
       Inches(0.45), Inches(0.10), Inches(12.0), Inches(0.54),
       Pt(20), True, P_WHITE, PP_ALIGN.LEFT)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.45), Inches(0.64), Inches(12.0), Inches(0.32),
           Pt(9.5), False, RGBColor(0xFF, 0xC0, 0xBF), PP_ALIGN.LEFT)


# =============================================================================
# SLIDE 1 — Cover
# =============================================================================
def slide_cover():
    sl = prs.slides.add_slide(COVER)
    for ph in sl.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "DSPy vs No-DSPy"
            ph.text_frame.paragraphs[0].runs[0].font.bold = True
        elif idx == 1:
            ph.text = "A Comparative Evaluation"
    return sl


# =============================================================================
# SLIDE 2 — Email + Flow Diagram + DSPy Intro
# =============================================================================
def slide_flow():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    hdr(sl,
        "Pipeline Architecture: 2-Agent Direct vs DSPy Orchestrator",
        "Real Japanese procurement email (Suzuki Sekkei RFQ) — processed by two competing pipelines")

    CY    = Inches(1.08)          # content top
    SPLIT = Inches(5.9)           # left column width
    GAP   = Inches(0.18)
    RX    = SPLIT + GAP           # right column start
    RW    = W - RX - Inches(0.25) # right column width

    # ── LEFT: email card ─────────────────────────────────────────────────────
    EH = Inches(6.28)
    box(sl, Inches(0.25), CY, SPLIT - Inches(0.25), EH, P_LGRAY, P_BORDER)
    box(sl, Inches(0.25), CY, SPLIT - Inches(0.25), Inches(0.32), P_DARK)
    tb(sl, "ACTUAL TEST EMAIL  (Japanese → English — Suzuki Sekkei RFQ)",
       Inches(0.38), CY + Inches(0.04), SPLIT - Inches(0.5), Inches(0.26),
       Pt(9), True, P_WHITE)

    lines = [
        ("From:",    "Kenta Yamamoto [山本健太]  <yamamoto.kenta@suzuki-sekkei.co.jp>"),
        ("To:",      "bids@globalparts.com"),
        ("Date:",    "May 26, 2026 — 10:32 AM JST"),
        ("Subject:", "[URGENT] Procurement — 6-Axis Arc Welding Robots"),
        ("",         "─" * 48),
        ("Item:",    "6-Axis Arc Welding Robot"),
        ("Model:",   "FANUC M-10iA/12 or Yaskawa Motoman MA1440 (equiv. ok)"),
        ("Qty:",     "18 units   (3-phase delivery: 6 + 6 + 6)"),
        ("Budget:",  "JPY 220,000,000  (excl. tax) — delivery & install incl."),
        ("",         ""),
        ("Phase 1:", "6 units  →  Sep 15, 2026"),
        ("Phase 2:", "6 units  →  Nov 30, 2026"),
        ("Phase 3:", "6 units  →  Jan 31, 2027"),
        ("Full Ops:","Mar 1, 2027"),
        ("",         "─" * 48),
        ("Specs:",   "IP67 | ±0.05 mm repeatability | 10 kg payload | 1,400 mm reach"),
        ("",         "EtherNet/IP controller required (PLC connection)"),
        ("AMC:",     "Min. 3 years — scope UNDEFINED in email"),
        ("Bid by:",  "Jun 20, 2026  17:00 JST"),
        ("Contact:", "Yamamoto Kenta [山本健太]  —  Head of Procurement"),
        ("Company:", "Suzuki Sekkei Co., Ltd. [鈴木設計株式会社]"),
        ("Address:", "3-18-1 Sakae, Naka-ku, Nagoya, Aichi 460-0008, Japan"),
    ]
    KW = Inches(0.78)
    VX = Inches(0.25) + KW + Inches(0.1)
    VW = SPLIT - KW - Inches(0.52)
    for i, (k, v) in enumerate(lines):
        y = CY + Inches(0.40) + i * Inches(0.252)
        if not v:
            continue
        if k:
            tb(sl, k, Inches(0.38), y, KW, Inches(0.24), Pt(7.5), True, P_DGRAY)
            tb(sl, v, VX, y, VW, Inches(0.24), Pt(7.5), False, P_DARK)
        else:
            tb(sl, v, Inches(0.38), y, SPLIT - Inches(0.5), Inches(0.22),
               Pt(7), False, P_MGRAY, PP_ALIGN.LEFT, True)

    # ── RIGHT TOP: DSPy intro ─────────────────────────────────────────────────
    INTRO_H = Inches(2.2)
    box(sl, RX, CY, RW, INTRO_H, P_PURP_LT, P_PURPLE)
    box(sl, RX, CY, RW, Inches(0.3), P_PURPLE)
    tb(sl, "Why DSPy?  (Declarative Self-improving Python)",
       RX + Inches(0.12), CY + Inches(0.04), RW - Inches(0.2), Inches(0.25),
       Pt(9.5), True, P_WHITE)

    dspy_pts = [
        ("Structured Pipelines",  "Chain LLM calls as composable, schema-typed modules — I/O is enforced"),
        ("Auditable Handoffs",    "Every step is inspectable via dspy.inspect_history(); no black-box calls"),
        ("Constrained Extraction","Signatures validate output structure — prevents schema drift and phantom fields"),
        ("Optimizable",           "BootstrapFewShot auto-tunes prompts over time without manual rewrites"),
    ]
    for i, (ttl, desc) in enumerate(dspy_pts):
        y = CY + Inches(0.38) + i * Inches(0.44)
        box(sl, RX + Inches(0.12), y + Inches(0.06), Inches(0.05), Inches(0.28), P_PURPLE)
        tb(sl, ttl + ":", RX + Inches(0.22), y + Inches(0.06),
           Inches(1.85), Inches(0.25), Pt(8.5), True, P_PURPLE)
        tb(sl, desc, RX + Inches(2.12), y + Inches(0.06),
           RW - Inches(2.28), Inches(0.25), Pt(8.5), False, P_DGRAY)

    # ── RIGHT BOTTOM: flow diagram ────────────────────────────────────────────
    FY = CY + INTRO_H + Inches(0.14)
    FH = EH - INTRO_H - Inches(0.14)
    box(sl, RX, FY, RW, FH, P_WHITE, P_BORDER)
    box(sl, RX, FY, RW, Inches(0.28), RGBColor(0x2D, 0x2D, 0x2D))
    tb(sl, "Pipeline Flow Comparison",
       RX + Inches(0.12), FY + Inches(0.04), RW - Inches(0.2), Inches(0.22),
       Pt(9), True, P_WHITE)

    def flow_row(base_y, label, label_color, nodes, node_color, node_line):
        tb(sl, label, RX + Inches(0.14), base_y,
           RW - Inches(0.2), Inches(0.22), Pt(8.5), True, label_color)
        BOX_Y  = base_y + Inches(0.26)
        BOX_H2 = Inches(0.66)
        AROW_Y = BOX_Y + Inches(0.28)
        x = RX + Inches(0.14)
        for j, (lbl1, lbl2, bw) in enumerate(nodes):
            box(sl, x, BOX_Y, bw, BOX_H2, P_LGRAY if lbl1 == "Email" else
                (P_PURP_LT if node_color == P_PURPLE else P_RED_LT if node_color == P_RED else P_GREEN_LT),
                P_BORDER if lbl1 == "Email" else node_line)
            tb(sl, lbl1, x + Inches(0.04), BOX_Y + Inches(0.06), bw - Inches(0.08),
               Inches(0.22), Pt(8), True,
               P_DGRAY if lbl1 == "Email" else node_color, PP_ALIGN.CENTER)
            if lbl2:
                tb(sl, lbl2, x + Inches(0.04), BOX_Y + Inches(0.32), bw - Inches(0.08),
                   Inches(0.28), Pt(7), False, P_DGRAY, PP_ALIGN.CENTER)
            x += bw
            if j < len(nodes) - 1:
                tb(sl, "→", x, AROW_Y - Inches(0.14),
                   Inches(0.26), Inches(0.32), Pt(14), True, P_MGRAY, PP_ALIGN.CENTER)
                x += Inches(0.26)
        # description under boxes
        DESC_Y = BOX_Y + BOX_H2 + Inches(0.04)
        x = RX + Inches(0.14)
        for j, (lbl1, lbl2, bw) in enumerate(nodes):
            if lbl2:
                tb(sl, lbl2.split("\n")[0][:28] if len(nodes) > 3 else "",
                   x, DESC_Y, bw, Inches(0.2), Pt(6.5), False, P_MGRAY, PP_ALIGN.CENTER)
            x += bw + Inches(0.26) if j < len(nodes) - 1 else 0

    # Case 1 — No-DSPy
    R1_Y = FY + Inches(0.38)
    tb(sl, "Case 1 — No DSPy  (2-Agent Direct):",
       RX + Inches(0.14), R1_Y, RW - Inches(0.2), Inches(0.22),
       Pt(8.5), True, P_RED)
    BOX1_Y = R1_Y + Inches(0.26)
    BH = Inches(0.70)
    AY = BOX1_Y + Inches(0.30)

    nodes_b = [
        ("Email", "Input", Inches(1.0)),
        ("Agent-1", "Extract\n+ Audit", Inches(2.2)),
        ("Agent-2", "Enrich\n+ Validate", Inches(2.2)),
        ("Output", "JSON", Inches(0.9)),
    ]
    x = RX + Inches(0.14)
    for j, (n1, n2, nw) in enumerate(nodes_b):
        fill = P_LGRAY if n1 in ("Email", "Output") else P_RED_LT
        border = P_BORDER if n1 in ("Email", "Output") else P_RED
        box(sl, x, BOX1_Y, nw, BH, fill, border)
        tb(sl, n1, x + Inches(0.04), BOX1_Y + Inches(0.08), nw - Inches(0.08), Inches(0.22),
           Pt(8.5), True, P_DGRAY if n1 in ("Email","Output") else P_RED, PP_ALIGN.CENTER)
        tb(sl, n2, x + Inches(0.04), BOX1_Y + Inches(0.34), nw - Inches(0.08), Inches(0.3),
           Pt(7), False, P_DGRAY, PP_ALIGN.CENTER)
        x += nw
        if j < len(nodes_b) - 1:
            tb(sl, "→", x, AY - Inches(0.10), Inches(0.28), Inches(0.30),
               Pt(13), True, P_MGRAY, PP_ALIGN.CENTER)
            x += Inches(0.28)

    # Case 2 — DSPy
    R2_Y = BOX1_Y + BH + Inches(0.55)
    tb(sl, "Case 2 — With DSPy  (Orchestrator):",
       RX + Inches(0.14), R2_Y, RW - Inches(0.2), Inches(0.22),
       Pt(8.5), True, P_GREEN)
    BOX2_Y = R2_Y + Inches(0.26)
    AY2 = BOX2_Y + Inches(0.30)

    nodes_a = [
        ("Email",             "Input",              Inches(1.0)),
        ("DSPy",              "Predict\n(Extract)",  Inches(1.65)),
        ("DSPy",              "Predict\n(Audit)",    Inches(1.65)),
        ("Agent",             "Enrich",              Inches(1.65)),
        ("Output",            "JSON",                Inches(0.9)),
    ]
    x = RX + Inches(0.14)
    for j, (n1, n2, nw) in enumerate(nodes_a):
        is_dspy  = n1 == "DSPy"
        is_agent = n1 == "Agent"
        fill   = P_LGRAY if n1 in ("Email","Output") else (P_PURP_LT if is_dspy else P_GREEN_LT)
        border = P_BORDER if n1 in ("Email","Output") else (P_PURPLE if is_dspy else P_GREEN)
        fcol   = P_DGRAY if n1 in ("Email","Output") else (P_PURPLE if is_dspy else P_GREEN)
        box(sl, x, BOX2_Y, nw, BH, fill, border)
        tb(sl, n1, x + Inches(0.04), BOX2_Y + Inches(0.08), nw - Inches(0.08), Inches(0.22),
           Pt(8.5), True, fcol, PP_ALIGN.CENTER)
        tb(sl, n2, x + Inches(0.04), BOX2_Y + Inches(0.34), nw - Inches(0.08), Inches(0.3),
           Pt(7), False, P_DGRAY, PP_ALIGN.CENTER)
        x += nw
        if j < len(nodes_a) - 1:
            tb(sl, "→", x, AY2 - Inches(0.10), Inches(0.26), Inches(0.30),
               Pt(13), True, P_MGRAY, PP_ALIGN.CENTER)
            x += Inches(0.26)

    # Key difference callout
    KD_Y = BOX2_Y + BH + Inches(0.14)
    box(sl, RX + Inches(0.14), KD_Y, RW - Inches(0.28), Inches(0.52), P_AMBER_LT, P_AMBER)
    tb(sl, "Key Difference:  DSPy owns structured extraction and audit via typed Signatures — "
          "prevents schema drift, phantom fields, and missing domain defaults",
       RX + Inches(0.26), KD_Y + Inches(0.09), RW - Inches(0.52), Inches(0.36),
       Pt(8.5), False, P_AMBER, PP_ALIGN.LEFT)


# =============================================================================
# SHARED: draw metrics table (header + rows + optional left-accent bars)
# =============================================================================
def draw_metrics_table(sl, T_X, T_Y, col_widths, headers, rows,
                       row_h=Inches(0.64), hdr_h=Inches(0.36),
                       score_good=False):
    """
    col_widths : list of widths [metric, score, well, not]
    score_good : True → score col uses P_GREEN; False → P_RED
    Returns bottom Y of table.
    """
    # Header row
    x = T_X
    for i, (hd, cw) in enumerate(zip(headers, col_widths)):
        box(sl, x, T_Y, cw, hdr_h, P_DARK)
        align = PP_ALIGN.CENTER if i == 1 else PP_ALIGN.LEFT
        tb(sl, hd, x + Inches(0.08), T_Y + Inches(0.07),
           cw - Inches(0.12), hdr_h - Inches(0.08), Pt(9.5), True, P_WHITE, align)
        x += cw

    # Data rows
    for r, row in enumerate(rows):
        ry = T_Y + hdr_h + r * row_h
        x  = T_X
        bg_clr = P_LGRAY if r % 2 == 0 else P_WHITE
        for c, (cell, cw) in enumerate(zip(row, col_widths)):
            box(sl, x, ry, cw, row_h, bg_clr, P_BORDER, Pt(0.3))
            if c == 0:  # Metric
                tb(sl, cell, x + Inches(0.1), ry + Inches(0.1),
                   cw - Inches(0.15), row_h - Inches(0.14), Pt(9.5), True, P_DARK)
            elif c == 1:  # Score
                sc = P_GREEN if score_good else P_RED
                tb(sl, cell, x + Inches(0.04), ry + Inches(0.1),
                   cw - Inches(0.08), row_h - Inches(0.14), Pt(19), True, sc, PP_ALIGN.CENTER)
            elif c == 2:  # Went well
                box(sl, x, ry, Inches(0.04), row_h, P_GREEN)
                tb(sl, cell, x + Inches(0.1), ry + Inches(0.07),
                   cw - Inches(0.14), row_h - Inches(0.1), Pt(8), False, P_DGRAY)
            elif c == 3:  # Did NOT
                box(sl, x, ry, Inches(0.04), row_h, P_RED)
                tb(sl, cell, x + Inches(0.1), ry + Inches(0.07),
                   cw - Inches(0.14), row_h - Inches(0.1), Pt(8), False, P_DGRAY)
            x += cw

    return T_Y + hdr_h + len(rows) * row_h


# =============================================================================
# SLIDE 3 — No-DSPy Metrics
# =============================================================================
def slide_nodspy():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    hdr(sl,
        "No-DSPy Pipeline (2-Agent Direct) — Evaluation Metrics",
        "Pipeline B — Japanese RFQ (Suzuki Sekkei) — direct agent chaining without structured orchestration")

    T_X = Inches(0.28)
    T_Y = Inches(1.10)
    T_W = W - Inches(0.56)
    CW  = [Inches(1.7), Inches(1.1), T_W - Inches(1.7) - Inches(1.1) - Inches(5.0), Inches(5.0)]
    CW[2] = (T_W - CW[0] - CW[1] - CW[3])

    rows = [
        ("F1 Score\n(Extraction\nAccuracy)",
         "0.83",
         "• All 40 core facts extracted from Japanese email\n"
         "• Correct: item, qty, budget, specs, 3 delivery phases\n"
         "• Compliance standards (JIS B 8433, ISO 10218-1) correct",
         "• 2 phantom sections created (responding_party, secondary_contact)\n"
         "• 12 spurious NEEDS_INPUT — only 2 are genuine gaps\n"
         "• Structural errors reduce precision below 0.87 benchmark"),

        ("Faithfulness\n(Source\nGrounding)",
         "0.95",
         "• Readable English dates: 'September 15, 2026' ✓\n"
         "• Japanese traceability brackets throughout:\n"
         "   'Kenta Yamamoto [山本健太]'\n"
         "• Captured 'equivalent models acceptable' explicitly",
         "• Vendor (responding_party) fields injected — not in source\n"
         "• Risk flag 'No supplier details' wrong for buyer-side RFQ\n"
         "• 1 genuine open item: equivalent model pre-approval unclear"),

        ("Output\nRelevancy\n(Confidence)",
         "64%",
         "• Correct domain: procurement / request_for_quote\n"
         "• All technical specs extracted: IP67, ±0.05 mm, EtherNet/IP\n"
         "• open_items returned as array — correct schema format",
         "• 14 NEEDS_INPUT (12 phantom — not real data gaps)\n"
         "• 0/7 domain-standard defaults applied: payment_terms blank\n"
         "• ERP record requires manual completion before use"),

        ("Output\nCorrectness\n(Field Quality)",
         "0.80",
         "• Technical specs fully and correctly extracted\n"
         "• Budget breakdown (inclusions/exclusions) correct\n"
         "• Delivery batches (3 × 6 units) accurately structured",
         "• 0 domain defaults: governing_law, JCAA, force_majeure all blank\n"
         "• Only 45 of 55 expected fields filled\n"
         "• Resolution log inflated by phantom field decisions"),

        ("Hallucination\nRate",
         "15%",
         "• Core factual values grounded in source email\n"
         "• All dates, amounts, specs match original Japanese\n"
         "• Resolution log honest — escalated unknowable items",
         "• Schema inflation: vendor sections created from nothing\n"
         "• secondary_contact (6 fields) — no basis in buyer RFQ\n"
         "• Wrong risk flag: 'No supplier details may delay RFQ'"),
    ]

    hdrs = ["Metric", "Score", "What Went Well", "What Did NOT"]
    bottom = draw_metrics_table(sl, T_X, T_Y, CW, hdrs, rows,
                                row_h=Inches(0.98), score_good=False)

    # Summary banner
    GAP  = Inches(0.1)
    SY   = bottom + GAP
    box(sl, T_X, SY, T_W, Inches(0.56), P_RED_LT, P_RED)
    box(sl, T_X, SY, Inches(0.1), Inches(0.56), P_RED)
    tb(sl, "Pipeline B Summary:  Factual extraction is good (Faithfulness 0.95), but schema inflation, "
          "missing domain fills, and wrong risk flags make the record incomplete for ERP use. "
          "Wins Faithfulness only; loses F1, Confidence, Hallucination, and Correctness.",
       T_X + Inches(0.2), SY + Inches(0.08), T_W - Inches(0.3), Inches(0.42),
       Pt(8.5), False, P_DARK, PP_ALIGN.LEFT)


# =============================================================================
# SLIDE 4 — DSPy Metrics + Final Summary Comparison Table
# =============================================================================
def slide_dspy_metrics():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    hdr(sl,
        "DSPy Orchestrator — Evaluation Metrics & Final Comparison",
        "Pipeline A — Same Japanese RFQ — typed Signatures and structured audit vs direct agent chaining")

    T_X = Inches(0.28)
    T_Y = Inches(1.10)
    T_W = W - Inches(0.56)
    CW  = [Inches(1.7), Inches(1.1), Inches(0), Inches(5.0)]
    CW[2] = T_W - CW[0] - CW[1] - CW[3]

    rows_dspy = [
        ("F1 Score\n(Extraction\nAccuracy)",
         "0.87",
         "• 55 fields filled — 10 more than 2-Agent Direct (45)\n"
         "• 7 domain defaults applied: Net 30, JCAA, force majeure…\n"
         "• Only 2 genuine NEEDS_INPUT: rfq_id and reference_number",
         "• Duplicate keys: delivery_schedule appears twice in output\n"
         "• open_items returned as object {} not array [] — schema error\n"
         "• amc_contract_duration duplicated in domain_specific_fields"),

        ("Faithfulness\n(Source\nGrounding)",
         "0.92",
         "• All 40 core facts extracted and verified\n"
         "• Japanese text fully translated — no residual characters\n"
         "• Zero phantom sections — no vendor fields invented",
         "• Dates in ISO format (2026-09-15) instead of 'Sep 15, 2026'\n"
         "• Inferred 'Global Parts Co., Ltd.' from email address alone\n"
         "• 3 real procurement risks not surfaced in risk_flags"),

        ("Output\nRelevancy\n(Confidence)",
         "79%",
         "• ERP-ready record — all standard domain fields filled\n"
         "• Laws of Japan, JCAA arbitration, bid_validity_days applied\n"
         "• +23% confidence vs Pipeline B (79% vs 64%)",
         "• risk_flags: [] — 3 critical risks missed completely\n"
         "• AMC scope, budget tax exclusion, no RFQ ref — all un-flagged\n"
         "• Completeness 'mostly_complete' — 2 NEEDS_INPUT remain"),

        ("Output\nCorrectness\n(Field Quality)",
         "0.88",
         "• AuditSignature validates every field vs source email\n"
         "• Domain fills: Net 30, JCAA, force_majeure, evaluation_criteria\n"
         "• Resolution log documents all merge/drop decisions",
         "• open_items schema error: {} not [] — downstream parse issue\n"
         "• Delivery dates ISO format — minor but affects faithfulness\n"
         "• Resolution log only 2 entries (4 in Pipeline B)"),

        ("Hallucination\nRate",
         "8%",
         "• AuditSignature eliminates phantom sections — zero inflation\n"
         "• No responding_party or secondary_contact invented\n"
         "• Hallucination nearly halved vs Pipeline B (8% vs 15%)",
         "• 'Global Parts Co., Ltd.' inferred from email domain — unverified\n"
         "• Empty risk_flags is a silent miss — not a hallucination, but a gap\n"
         "• Both pipelines share this risk blind spot — top fix priority"),
    ]

    hdrs = ["Metric", "Score", "What Went Well", "What Did NOT"]
    bottom = draw_metrics_table(sl, T_X, T_Y, CW, hdrs, rows_dspy,
                                row_h=Inches(0.64), score_good=True)

    # ── FINAL SUMMARY TABLE ───────────────────────────────────────────────────
    SY  = bottom + Inches(0.12)
    SRH = Inches(0.30)  # summary row height
    SCW = [Inches(1.7), Inches(1.05), Inches(1.05), Inches(0.90), T_W - Inches(1.7) - Inches(1.05)*2 - Inches(0.90)]

    # Dark header banner
    box(sl, T_X, SY, T_W, Inches(0.28), P_DARK)
    tb(sl, "FINAL COMPARISON — Pipeline A (DSPy) vs Pipeline B (2-Agent Direct)",
       T_X + Inches(0.12), SY + Inches(0.04),
       T_W - Inches(0.2), Inches(0.22), Pt(9), True, P_WHITE)

    # Column headers
    CHY = SY + Inches(0.28)
    CHDRS = ["Metric", "No-DSPy", "DSPy", "Delta", "Finding"]
    x = T_X
    for i, (ch, cw) in enumerate(zip(CHDRS, SCW)):
        box(sl, x, CHY, cw, SRH, RGBColor(0x33, 0x33, 0x33))
        align = PP_ALIGN.CENTER if 0 < i < 4 else PP_ALIGN.LEFT
        tb(sl, ch, x + Inches(0.06), CHY + Inches(0.06),
           cw - Inches(0.1), SRH - Inches(0.08), Pt(8.5), True, P_WHITE, align)
        x += cw

    summary_rows = [
        ("F1 Score",         "0.83",  "0.87",  "+5%",   True,  "DSPy — 7 domain defaults; fewer structural errors"),
        ("Faithfulness",     "0.95",  "0.92",  "−3 pp", False, "No-DSPy — readable dates, Japanese traceability brackets"),
        ("Output Relevancy", "64%",   "79%",   "+23%",  True,  "DSPy — 55 vs 45 fields; only 2 genuine NEEDS_INPUT"),
        ("Hallucination",    "15%",   "8%",    "−47%",  True,  "DSPy — zero phantom sections; no schema inflation"),
        ("Output Correct.",  "0.80",  "0.88",  "+10%",  True,  "DSPy — AuditSignature prevents drift; domain fills applied"),
    ]

    for r, (metric, vb, va, delta, dspy_w, finding) in enumerate(summary_rows):
        ry  = CHY + SRH + r * SRH
        rbg = P_LGRAY if r % 2 == 0 else P_WHITE
        x   = T_X
        dc  = P_GREEN if dspy_w else P_AMBER
        cols = [metric, vb, va, delta, finding]
        for c, (cell, cw) in enumerate(zip(cols, SCW)):
            box(sl, x, ry, cw, SRH, rbg, P_BORDER, Pt(0.25))
            if c == 0:
                tb(sl, cell, x + Inches(0.06), ry + Inches(0.06),
                   cw - Inches(0.1), SRH - Inches(0.08), Pt(8.5), True, P_DARK)
            elif c == 1:
                tb(sl, cell, x + Inches(0.06), ry + Inches(0.06),
                   cw - Inches(0.1), SRH - Inches(0.08), Pt(9), True, P_RED, PP_ALIGN.CENTER)
            elif c == 2:
                tb(sl, cell, x + Inches(0.06), ry + Inches(0.06),
                   cw - Inches(0.1), SRH - Inches(0.08), Pt(9), True,
                   P_GREEN if dspy_w else P_AMBER, PP_ALIGN.CENTER)
            elif c == 3:
                box(sl, x + Inches(0.05), ry + Inches(0.05),
                    cw - Inches(0.1), SRH - Inches(0.1), dc)
                tb(sl, cell, x + Inches(0.05), ry + Inches(0.05),
                   cw - Inches(0.1), SRH - Inches(0.1), Pt(8), True, P_WHITE, PP_ALIGN.CENTER)
            else:
                tb(sl, cell, x + Inches(0.08), ry + Inches(0.06),
                   cw - Inches(0.14), SRH - Inches(0.08), Pt(8), False, P_DGRAY)
            x += cw

    # Verdict
    VY = CHY + SRH + len(summary_rows) * SRH + Inches(0.06)
    box(sl, T_X, VY, T_W, Inches(0.36), P_GREEN_LT, P_GREEN)
    box(sl, T_X, VY, Inches(0.1), Inches(0.36), P_GREEN)
    tb(sl, "VERDICT: DSPy Orchestrator wins 4 of 5 metrics — lower hallucination (8% vs 15%), "
          "better coverage (55 vs 45 fields), higher confidence (79% vs 64%), better correctness. "
          "No-DSPy wins Faithfulness (0.95 vs 0.92) via readable dates and Japanese brackets. "
          "Critical gap for BOTH: add RiskSignature to surface AMC scope, tax, and RFQ ID risks.",
       T_X + Inches(0.2), VY + Inches(0.05), T_W - Inches(0.3), Inches(0.3),
       Pt(8), False, P_DARK, PP_ALIGN.LEFT)


# ── Build ─────────────────────────────────────────────────────────────────────
slide_cover()
slide_flow()
slide_nodspy()
slide_dspy_metrics()

OUT = "dspy_evaluation_ver2.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
